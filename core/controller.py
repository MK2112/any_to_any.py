import os
import re
import fitz
import docx
import pptx
import shutil
import PyPDF2
import mammoth
import logging
import subprocess
import numpy as np
import utils.language_support as lang
from PIL import Image
from tqdm import tqdm
from pathlib import Path
from weasyprint import HTML
from utils.category import Category
from utils.prog_logger import ProgLogger
from core.audio_converter import AudioConverter
from core.movie_converter import MovieConverter
from core.utils.file_handler import FileHandler
from core.utils.misc import end_with_msg, gifs_to_frames
from core.utils.directory_watcher import DirectoryWatcher
from core.doc_converter import office_to_frames, DocumentConverter
from moviepy import (
    AudioFileClip,
    VideoFileClip,
    VideoClip,
    ImageClip,
    concatenate_videoclips,
    concatenate_audioclips,
    clips_array,
)

# TODO: Move to_markdown, to_pdf, to_subtitles, to_office to doc_converter
# TODO: Finalize doc_converter
# TODO: Move to_frames, to_bmp, to_webp, to_gif to image_converter
# TODO: Finalize image_converter
# TODO: Add converter-wise tests


class Controller:
    """
    Taking an input directory of files, convert them to a multitude of formats.
    Interact with the script using the command line arguments or the web interface.
    Run via any_to_any.py script.
    """

    def __init__(self):
        # Setting up progress logger
        self.prog_logger = ProgLogger()
        # Get locale
        self.locale = lang.get_system_language()
        # Setting up event logger and format
        logging.basicConfig(
            level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s"
        )
        self.event_logger = logging.getLogger(__name__)
        self.file_handler = FileHandler(self.event_logger, self.locale)
        # Setup distinct converters
        self.audio_converter = AudioConverter(
            self.file_handler, self.prog_logger, self.event_logger, self.locale
        )
        self.movie_converter = MovieConverter(
            self.file_handler, self.prog_logger, self.event_logger, self.locale
        )
        self.doc_converter = DocumentConverter(
            self.file_handler, self.prog_logger, self.event_logger, self.locale
        )
        # Setting up a dictionary of supported formats and respective information
        self._supported_formats = {
            Category.AUDIO: {
                "mp3": "libmp3lame",
                "flac": "flac",
                "aac": "aac",
                "ac3": "ac3",
                "dts": "dts",
                "ogg": "libvorbis",
                "wma": "wmav2",
                "wav": "pcm_s16le",
                "m4a": "aac",
                "aiff": "pcm_s16le",
                "weba": "libopus",
                "mka": "libvorbis",
                "wv": "wavpack",
                "tta": "tta",
                "m4b": "aac",
                "eac3": "eac3",
                "spx": "libvorbis",
                "mp2": "mp2",
                "caf": "pcm_s16be",
                "au": "pcm_s16be",
                "oga": "libvorbis",
                "opus": "libopus",
                "m3u8": "pcm_s16le",
                "w64": "pcm_s16le",
                "mlp": "mlp",
                "adts": "aac",
                "sbc": "sbc",
                "thd": "truehd",
            },
            Category.IMAGE: {
                "gif": self.to_gif,
                "png": self.to_frames,
                "jpeg": self.to_frames,
                "jpg": self.to_frames,
                "bmp": self.to_bmp,
                "webp": self.to_webp,
                "tiff": self.to_frames,
                "tga": self.to_frames,
                "ps": self.to_frames,
                "ico": self.to_frames,
                "eps": self.to_frames,
                "jpeg2000": self.to_frames,
                "im": self.to_frames,
                "pcx": self.to_frames,
                "ppm": self.to_frames,
            },
            Category.DOCUMENT: {
                "md": self.doc_converter.to_markdown,
                "pdf": self.to_pdf,
                "docx": self.to_office,
                "pptx": self.to_office,
                "srt": self.to_subtitles,
            },
            Category.MOVIE: {
                "webm": "libvpx",
                "mov": "libx264",
                "mkv": "libx264",
                "avi": "libx264",
                "mp4": "libx264",
                "wmv": "wmv2",
                "flv": "libx264",
                "mjpeg": "mjpeg",
                "m2ts": "mpeg2video",
                "3gp": "libx264",
                "3g2": "libx264",
                "asf": "wmv2",
                "vob": "mpeg2video",
                "ts": "hevc",
                "raw": "rawvideo",
                "mpg": "mpeg2video",
                "mxf": "mpeg2video",
                "drc": "libx265",
                "swf": "flv",
                "f4v": "libx264",
                "m4v": "libx264",
                "mts": "mpeg2video",
                "m2v": "mpeg2video",
                "yuv": "rawvideo",
            },
            Category.MOVIE_CODECS: {
                "av1": ["libaom-av1", "mkv"],  # [lib, fallback]
                "avc": ["libx264", "mp4"],
                "vp9": ["libvpx-vp9", "mp4"],
                "h265": ["libx265", "mkv"],
                "h264": ["libx264", "mkv"],
                "h263p": ["h263p", "mkv"],
                "xvid": ["libxvid", "mp4"],
                "mpeg4": ["mpeg4", "mp4"],
                "theora": ["libtheora", "ogv"],
                "mpeg2": ["mpeg2video", "mp4"],
                "mpeg1": ["mpeg1video", "mp4"],
                "hevc": ["libx265", "mkv"],
                "prores": ["prores", "mkv"],
                "vp8": ["libvpx", "webm"],
                "huffyuv": ["huffyuv", "mkv"],
                "ffv1": ["ffv1", "mkv"],
                "ffvhuff": ["ffvhuff", "mkv"],
                "v210": ["v210", "mkv"],
                "v410": ["v410", "mkv"],
                "v308": ["v308", "mkv"],
                "v408": ["v408", "mkv"],
                "zlib": ["zlib", "mkv"],
                "qtrle": ["qtrle", "mkv"],
                "snow": ["snow", "mkv"],
                "svq1": ["svq1", "mkv"],
                "utvideo": ["utvideo", "mkv"],
                "cinepak": ["cinepak", "mkv"],
                "msmpeg4": ["msmpeg4", "mkv"],
                "h264_nvenc": ["h264_nvenc", "mp4"],
                "vpx": ["libvpx", "webm"],
                "h264_rgb": ["libx264rgb", "mkv"],
                "mpeg2video": ["mpeg2video", "mpg"],
                "prores_ks": ["prores_ks", "mkv"],
                "vc2": ["vc2", "mkv"],
                "flv1": ["flv", "flv"],
            },
            Category.PROTOCOLS: {
                "hls": ["hls", "mkv"],
                "dash": ["dash", "mkv"],
            },
        }

        # Used in CLI information output
        self.supported_formats = [
            format
            for formats in self._supported_formats.values()
            for format in formats.keys()
        ]

        # Indicates if the script is being run from the web interface
        self.web_flag = False
        # Host address for the web interface
        self.web_host = None

    def _audio_bitrate(self, format: str, quality: str) -> str:
        # Return bitrate for audio conversion
        # If formats allow for a higher bitrate, we shift our scale accordingly
        if format in [
            "flac",
            "wav",
            "aac",
            "aiff",
            "eac3",
            "dts",
            "au",
            "wv",
            "tta",
            "mlp",
        ]:
            return {
                "high": "500k",
                "medium": "320k",
                "low": "192k",
            }.get(quality, None)
        else:
            return {
                "high": "320k",
                "medium": "192k",
                "low": "128k",
            }.get(quality, None)

    def run(
        self,
        input_path_args: list,
        format: str,
        output: str,
        framerate: int,
        quality: str,
        merge: bool,
        concat: bool,
        delete: bool,
        across: bool,
        recursive: bool,
        dropzone: bool,
        language: str,
    ) -> None:
        # Main function, convert media files to defined formats
        # or merge or concatenate, according to the arguments
        input_paths = []
        input_path_args = (
            input_path_args
            if input_path_args is not None
            else [os.path.dirname(os.getcwd())]
        )

        self.across = across
        self.recursive = recursive

        # User-set Language
        if language is not None:
            # we expect a language code like "en_US" or "pl_PL"
            if re.match(r"^[a-z]{2}_[A-Z]{2}$", language) and language in list(
                lang.LANGUAGE_CODES.keys()
            ):
                self.locale = lang.LANGUAGE_CODES[language]
            else:
                self.event_logger.warning(
                    f"[!] {lang.get_translation('error', self.locale)}: {lang.get_translation('lang_not_supported', self.locale)}"
                )

        for _, arg in enumerate(input_path_args):
            # Custom handling of multiple input paths
            # (e.g. "-1 path1 -2 path2 -n pathn")
            if arg.startswith("-") and arg[1:].isdigit():
                input_paths.append(arg[2:])
            else:
                try:
                    input_paths[-1] = (input_paths[-1] + f" {arg}").strip()
                except IndexError:
                    input_paths.append(arg)

        if len(input_paths) == 1:
            self.output = output if output is not None else input_paths[0]
        else:
            # len(input_paths) > 1
            if not across:
                self.output = output if output is not None else None
            else:
                self.output = (
                    output if output is not None else os.path.dirname(os.getcwd())
                )

        # Check if the output dir exists - Create it otherwise
        if self.output is not None and not os.path.exists(self.output):
            os.makedirs(self.output)

        # No format means no conversion (but maybe merge || concat)
        self.target_format = format.lower() if format is not None else None
        self.framerate = framerate  # Possibly no framerate means same as input
        self.delete = delete  # Delete mp4 files after conversion
        # Check if quality is set, if not, set it to None
        self.quality = (
            (quality.lower() if quality.lower() in ["high", "medium", "low"] else None)
            if quality is not None
            else None
        )

        # Merge movie files with equally named audio files
        self.merging = merge
        # Concatenate files of same type (img/movie/audio) back to back
        self.concatenating = concat

        file_paths = {}
        was_none, found_files = False, False

        # Check if input paths are given
        for input_path in input_paths:
            input_path = os.path.abspath(input_path)
            try:
                file_paths = self.file_handler.get_file_paths(
                    input_path, file_paths, self._supported_formats
                )
            except FileNotFoundError:
                end_with_msg(
                    self.event_logger,
                    FileNotFoundError,
                    f"[!] {lang.get_translation('error', self.locale)}: {lang.get_translation('no_dir_exist', self.locale).replace('[dir]', f'{input_path}')}",
                )
            # Make self.input hold a directory
            if os.path.isfile(str(input_path)):
                self.input = os.path.dirname(input_path)
            else:
                self.input = input_path

            # If no output path set or derived by the script by now, throw an error
            if not self.output:
                end_with_msg(
                    self.event_logger,
                    ValueError,
                    f"[!] {lang.get_translation('error', self.locale)}: {lang.get_translation('no_out_multi_in', self.locale)}",
                )

            # If output is just a file for whatever reason, turn it into directory
            if os.path.isfile(self.output):
                self.output = os.path.dirname(self.output)

            # Unify path formatting, still works like any other string
            self.input, self.output = Path(self.input), Path(self.output)

            # Cut setup for media conversion short, begin dropzone mode
            if dropzone:
                if self.input == self.output or self.input in self.output.parents:
                    end_with_msg(
                        self.event_logger,
                        None,
                        f"[!] {lang.get_translation('error', self.locale)}: {lang.get_translation('dropzone_diff', self.locale)}",
                    )
                self.event_logger.info(
                    f"[+] {lang.get_translation('dropzone_active', self.locale)} {self.input}"
                )
                self.watchdropzone(self.input)
                return

            # What if that directory contains subdirectories?
            if self.recursive:
                # Check if there is any subdirectory in the input path
                if any(entry.is_dir() for entry in os.scandir(self.input)):
                    file_paths = {}  # Reset, we go through everything again to be sure
                    # If the user set the recursive option, also scan every non-empty subdirectory
                    for root, _, files in os.walk(self.input):
                        # root should be a directory with >=1 file to be considered
                        try:
                            file_paths = self.file_handler.get_file_paths(
                                root, file_paths, self._supported_formats
                            )
                        except FileNotFoundError:
                            end_with_msg(
                                self.event_logger,
                                FileNotFoundError,
                                f"[!] {lang.get_translation('error', self.locale)}: {lang.get_translation('no_dir_exist', self.locale).replace('[dir]', f'{root}')}",
                            )

            if not any(file_paths.values()):
                if len(input_paths) > 1:
                    self.event_logger.info(
                        f"[!] {lang.get_translation('no_media_found', self.locale).replace('[path]', f"'{str(input_path)}'")} - {lang.get_translation('skipping', self.locale)}"
                    )
                else:
                    end_with_msg(
                        self.event_logger,
                        None,
                        f"[!] {lang.get_translation('no_media_found', self.locale).replace('[path]', f"'{str(input_path)}'")}",
                    )
                continue

            # If no output given, output is set to the input path
            if not across:
                if self.output is None or was_none:
                    self.output = os.path.dirname(input_path)
                    was_none = True
                self.process_file_paths(file_paths)
                found_files = len(file_paths) > 0 if not found_files else found_files
                file_paths = {}

        # If multiple input paths are given, yet no output, output is set to the first input path
        if across:
            if self.output is None:
                self.output = os.path.dirname(input_paths[0])
            self.process_file_paths(file_paths)

        # Check if file_paths is empty
        if across and len(file_paths) == 0 or not found_files:
            self.event_logger.warning(
                f"{lang.get_translation('no_media_general', self.locale)}"
            )
        self.event_logger.info(
            f"[+] {lang.get_translation('job_finished', self.locale)}"
        )

    def process_file_paths(self, file_paths: dict) -> None:
        # Check if value associated to format is tuple/string or function to call specific conversion
        if self.target_format in self._supported_formats[Category.MOVIE].keys():
            self.movie_converter.to_movie(
                input=self.input,
                output=self.output,
                recursive=self.recursive,
                file_paths=file_paths,
                format=self.target_format,
                framerate=self.framerate,
                codec=self._supported_formats[Category.MOVIE][self.target_format],
                delete=self.delete,
            )
        elif self.target_format in self._supported_formats[Category.AUDIO].keys():
            self.audio_converter.to_audio(
                file_paths=file_paths,
                format=self.target_format,
                codec=self._supported_formats[Category.AUDIO][self.target_format],
                recursive=self.recursive,
                bitrate=self._audio_bitrate(self.target_format, self.quality),
                input=self.input,
                output=self.output,
                delete=self.delete,
            )
        elif (
            self.target_format in self._supported_formats[Category.MOVIE_CODECS].keys()
        ):
            self.movie_converter.to_codec(
                input=self.input,
                output=self.output,
                format=self.target_format,
                recursive=self.recursive,
                file_paths=file_paths,
                framerate=self.framerate,
                codec=self._supported_formats[Category.MOVIE_CODECS][
                    self.target_format
                ],
                delete=self.delete,
            )
        elif self.target_format in self._supported_formats[Category.IMAGE].keys():
            self._supported_formats[Category.IMAGE][self.target_format](
                file_paths, self.target_format
            )
        elif self.target_format in self._supported_formats[Category.DOCUMENT].keys():
            self._supported_formats[Category.DOCUMENT][self.target_format](self.input, self.output, file_paths, self.target_format, self.delete)
        elif self.target_format in self._supported_formats[Category.PROTOCOLS].keys():
            self.movie_converter.to_protocol(
                output=self.output,
                file_paths=file_paths,
                supported_formats=self._supported_formats,
                protocol=self._supported_formats[Category.PROTOCOLS][
                    self.target_format
                ],
                delete=self.delete,
            )
        elif self.merging:
            self.merge(file_paths, getattr(self, "across", False))
        elif self.concatenating:
            self.concat(file_paths, self.target_format)
        else:
            # Handle unsupported formats here
            end_with_msg(
                self.event_logger,
                ValueError,
                f"[!] {lang.get_translation('error', self.locale)}: {lang.get_translation('output_list', self.locale).replace('[list]', str(list(self.supported_formats)))}",
            )

    def watchdropzone(self, watch_path: str) -> None:
        # Watch a directory for new files and process them automatically
        def handle_file_event(event_type: str, file_path: str) -> None:
            if event_type == "created":
                try:
                    self.event_logger.info(
                        f"[>] {lang.get_translation('dropzone_new_file', self.locale)}: {file_path}"
                    )
                    if os.path.isfile(file_path):
                        # Create a temporary instance with distinct settings
                        temp_instance = self.__class__()
                        # Configure the instance with current settings
                        temp_instance.output = self.output
                        temp_instance.delete = (
                            True  # Delete original files after processing
                        )
                        temp_instance.locale = self.locale
                        temp_instance.framerate = self.framerate
                        temp_instance.quality = self.quality
                        temp_instance.merging = self.merging
                        temp_instance.concatenating = self.concatenating
                        temp_instance.recursive = True
                        temp_instance.event_logger = self.event_logger
                        temp_instance.file_handler = self.file_handler

                        # Process the file
                        file_paths = temp_instance.file_handler.get_file_paths(
                            [file_path]
                        )
                        if file_paths:
                            try:
                                temp_instance.process_files(file_paths)
                            except Exception as e:
                                self.event_logger.error(
                                    f"{lang.get_translation('error', self.locale)}: {file_path} - {str(e)}"
                                )
                except Exception as e:
                    self.event_logger.error(
                        f"{lang.get_translation('error', self.locale)}: {file_path} - {str(e)}"
                    )

        try:
            # Validate watch path
            watch_path = os.path.abspath(watch_path)
            if not os.path.exists(watch_path):
                self.event_logger.error(
                    f"{lang.get_translation('not_exist_not_dir', self.locale)}: {watch_path}"
                )
                return
            if not os.path.isdir(watch_path):
                self.event_logger.error(
                    f"{lang.get_translation('watch_not_dir', self.locale)}: {watch_path}"
                )
                return

            self.event_logger.info(
                f"[>] {lang.get_translation('watch_start', self.locale)}: {watch_path}"
            )
            self.event_logger.info(
                f"[>] {lang.get_translation('watch_stop', self.locale)}"
            )

            # Start the directory watcher with error handling
            with DirectoryWatcher(watch_path, handle_file_event) as watcher:
                watcher.watch()

        except KeyboardInterrupt:
            self.event_logger.info(
                f"\n[!] {lang.get_translation('watch_halt', self.locale)}"
            )
        except Exception as e:
            self.event_logger.error(
                f"{lang.get_translation('error', self.locale)}: {str(e)}"
            )
            raise

    def to_subtitles(self, file_paths: dict, format: str) -> None:
        for movie_path_set in file_paths[Category.MOVIE]:
            input_path = self.file_handler.join_back(movie_path_set)
            out_path = os.path.abspath(
                os.path.join(self.output, f"{movie_path_set[1]}.srt")
            )
            self.event_logger.info(
                f"[+] {lang.get_translation('extract_subtitles', self.locale)} '{input_path}'"
            )
            try:
                # Use FFmpeg to extract subtitles
                _ = subprocess.run(
                    [
                        "ffmpeg",
                        "-i",
                        input_path,
                        "-map",
                        "0:s:0",  # Selects first subtitle stream
                        "-c:s",
                        "srt",
                        out_path,
                    ],
                    capture_output=True,
                    text=True,
                )

                if os.path.exists(out_path) and os.path.getsize(out_path) > 0:
                    self.event_logger.info(
                        f"[+] {lang.get_translation('subtitles_success', self.locale)} '{out_path}'"
                    )
                    self.file_handler.post_process(
                        movie_path_set, out_path, self.delete, show_status=False
                    )
                else:
                    # Try extracting closed captions when direct extract fails (found mostly in MP4 and MKV)
                    self.event_logger.info(
                        f"[!] {lang.get_translation('extract_subtitles_alt', self.locale)}"
                    )
                    _ = subprocess.run(
                        ["ffmpeg", "-i", input_path, "-c:s", format, out_path],
                        capture_output=True,
                        text=True,
                    )
                    if os.path.exists(out_path) and os.path.getsize(out_path) > 0:
                        self.event_logger.info(
                            f"[+] {lang.get_translation('embed_subtitles_success')} '{out_path}'"
                        )
                        self.file_handler.post_process(
                            movie_path_set, out_path, self.delete, show_status=False
                        )
                    else:
                        self.event_logger.info(
                            f"[!] {lang.get_translation('embed_subtitles_fail', self.locale)} '{input_path}'"
                        )
            except Exception as e:
                self.event_logger.info(
                    f"[!] {lang.get_translation('extract_subtitles_fail', self.locale)} {str(e)}"
                )
                try:
                    subprocess.run(["ffmpeg", "-version"], capture_output=True)
                except FileNotFoundError:
                    self.event_logger.info(
                        f"[!] {lang.get_translation('ffmpeg_not_found', self.locale)}"
                    )
                    break

    def to_office(self, file_paths: dict, format: str) -> None:
        def _new_container():
            return pptx.Presentation() if format == "pptx" else docx.Document()

        def _add_page(container):
            return (
                container.slides.add_slide(container.slide_layouts[5])
                if format == "pptx"
                else container
            )

        def _place_img(page, img_path, full_page=False):
            with Image.open(img_path) as im:
                w, h = im.size
            if format == "pptx":
                max_h, max_w = container.slide_height, container.slide_width
                if full_page or h > max_h or w > max_w:
                    if h / max_h >= w / max_w:
                        page.shapes.add_picture(img_path, 0, 0, height=max_h)
                    else:
                        page.shapes.add_picture(img_path, 0, 0, width=max_w)
                else:
                    page.shapes.add_picture(img_path, 0, 0)
            elif format == "docx":
                try:
                    page.add_picture(img_path, height=docx.shared.Inches(7))
                except docx.image.exceptions.UnexpectedEndOfFileError as e:
                    self.event_logger.info(e)

        gifs_to_frames(self.output, file_paths, self.file_handler)

        for image_path_set in file_paths[Category.IMAGE]:
            out_path = os.path.abspath(
                os.path.join(self.output, f"{image_path_set[1]}.{format}")
            )
            container = _new_container()
            if image_path_set[2] == "gif":
                frame_dir = os.path.join(self.output, image_path_set[1])
                for frame in tqdm(sorted(os.listdir(frame_dir))):
                    if not frame.endswith(".png"):
                        continue
                    page = _add_page(container)
                    _place_img(page, os.path.join(frame_dir, frame), full_page=True)
                    if format == "docx":
                        page.add_paragraph(f"Image: {frame}")
                        page.add_page_break()
                shutil.rmtree(frame_dir)
            else:
                page = _add_page(container)
                _place_img(
                    page, self.file_handler.join_back(image_path_set), full_page=True
                )
                if format == "docx":
                    page.add_paragraph(f"Image: {image_path_set[1]}")
            container.save(out_path)
            self.file_handler.post_process(image_path_set, out_path, self.delete)

        for movie_path_set in file_paths[Category.MOVIE]:
            if not self.file_handler.has_visuals(movie_path_set):
                continue

            out_path = os.path.abspath(
                os.path.join(self.output, f"{movie_path_set[1]}.{format}")
            )
            container = _new_container()
            clip = VideoFileClip(
                self.file_handler.join_back(movie_path_set),
                audio=False,
                fps_source="tbr",
            )
            digits = len(str(int(clip.duration * clip.fps)))

            for idx, frame in tqdm(
                enumerate(clip.iter_frames(fps=clip.fps, dtype="uint8"))
            ):
                frame_png = os.path.join(
                    self.output, f"{movie_path_set[1]}-temp-{idx:0{digits}d}.png"
                )
                Image.fromarray(frame).save(frame_png, format="PNG")
                page = _add_page(container)
                _place_img(page, frame_png, full_page=format == "pptx")
                if format == "docx":
                    page.add_paragraph(f"Frame: {idx}")
                    page.add_page_break()
                os.remove(frame_png)

            clip.close()
            container.save(out_path)
            self.file_handler.post_process(movie_path_set, out_path, self.delete)

        if format == "docx":
            for document_path_set in file_paths[Category.DOCUMENT]:
                if document_path_set[2] == format:
                    continue

                out_path = os.path.abspath(
                    os.path.join(self.output, f"{document_path_set[1]}.docx")
                )
                doc = docx.Document()

                if document_path_set[2] == "pptx":
                    pptx_path = self.file_handler.join_back(document_path_set)
                    presentation = pptx.Presentation(pptx_path)
                    for slide in tqdm(
                        presentation.slides, desc=f"{document_path_set[1]}"
                    ):
                        if slide.shapes.title.text != "":
                            # Add title to the document
                            doc.add_paragraph(
                                f"Slide {slide.slide_id} Title: {slide.shapes.title.text}"
                            )
                        for shape in slide.shapes:
                            if hasattr(shape, "image"):
                                img = shape.image
                                slide_img = os.path.join(
                                    self.output,
                                    f"{document_path_set[1]}-slide{slide.slide_id}.{img.ext}",
                                )
                                with open(slide_img, "wb") as f:
                                    f.write(img.blob)
                                doc.add_paragraph(f"Slide {slide.slide_id} Image:")
                                doc.add_picture(slide_img, width=docx.shared.Inches(5))
                                os.remove(slide_img)
                            if hasattr(shape, "text"):
                                if shape.text != "":
                                    doc.add_paragraph(
                                        f"Slide {slide.slide_id} Text: {shape.text}"
                                    )
                        doc.add_page_break()
                elif document_path_set[2] == "pdf":
                    pdf = fitz.open(self.file_handler.join_back(document_path_set))
                    for pnum in tqdm(range(len(pdf))):
                        page = pdf.load_page(pnum)
                        txt = page.get_text().strip()
                        if txt:
                            doc.add_paragraph(f"Page {pnum + 1} Text:")
                            for line in txt.splitlines():
                                doc.add_paragraph(line)
                        imgs = page.get_images(full=True)
                        if imgs:
                            doc.add_paragraph(f"Page {pnum + 1} Images:")
                        for i, img in enumerate(imgs):
                            xref = img[0]
                            pix = pdf.extract_image(xref)
                            img_bytes, img_ext = pix["image"], pix["ext"]
                            tmp_img = os.path.join(
                                self.output,
                                f"{document_path_set[1]}-page{pnum + 1}-{i}.{img_ext}",
                            )
                            with open(tmp_img, "wb") as f:
                                f.write(img_bytes)
                            doc.add_picture(tmp_img, width=docx.shared.Inches(5))
                            os.remove(tmp_img)
                        doc.add_page_break()
                doc.save(out_path)
                self.file_handler.post_process(document_path_set, out_path, self.delete)

    def to_pdf(self, file_paths: dict, format: str) -> None:
        # Convert GIFs to Frames using to_frames
        # Produces a folder with gif frame for each gif
        self._gifs_to_frames(file_paths)
        # Convert Images to PDF
        for image_path_set in file_paths[Category.IMAGE]:
            # Convert image to pdf
            pdf_path = ""
            if image_path_set[2] != "gif":
                doc = fitz.open()
                img = fitz.Pixmap(self.file_handler.join_back(image_path_set))
                rect = fitz.Rect(0, 0, img.width, img.height)
                page = doc.new_page(width=rect.width, height=rect.height)
                page.insert_image(rect, pixmap=img)
                pdf_path = os.path.abspath(
                    os.path.join(self.output, f"{image_path_set[1]}.{format}")
                )
                doc.save(pdf_path)
                doc.close()
                self.file_handler.post_process(image_path_set, pdf_path, self.delete)
            elif image_path_set[2] == "gif":
                # We suppose the gif was converted to frames and we have a folder of pngs
                # All pngs shall be merged into one pdf
                gif_frame_path = os.path.join(self.output, image_path_set[1])
                pdf_path = os.path.abspath(
                    os.path.join(self.output, f"{image_path_set[1]}.{format}")
                )
                doc = fitz.open()
                for frame in sorted(os.listdir(gif_frame_path)):
                    if frame.endswith(".png"):
                        img = fitz.Pixmap(os.path.join(gif_frame_path, frame))
                        rect = fitz.Rect(0, 0, img.width, img.height)
                        page = doc.new_page(width=rect.width, height=rect.height)
                        page.insert_image(rect, pixmap=img)
                doc.save(pdf_path)
                doc.close()
                # Remove the gif frame folder
                shutil.rmtree(gif_frame_path)
                self.file_handler.post_process(image_path_set, pdf_path, self.delete)
        # Convert Movies to PDF
        for movie_path_set in file_paths[Category.MOVIE]:
            if self.file_handler.has_visuals(movie_path_set):
                clip = VideoFileClip(
                    self.file_handler.join_back(movie_path_set),
                    audio=False,
                    fps_source="tbr",
                )
                pdf_path = os.path.abspath(
                    os.path.join(self.output, f"{movie_path_set[1]}.{format}")
                )
                num_digits = len(str(int(clip.duration * clip.fps)))
                doc = fitz.open()
                for i, frame in tqdm(
                    enumerate(clip.iter_frames(fps=clip.fps, dtype="uint8")),
                ):
                    frame_path = os.path.abspath(
                        os.path.join(
                            self.output,
                            f"{movie_path_set[1]}-temp-{i:0{num_digits}d}.png",
                        )
                    )
                    Image.fromarray(frame).save(frame_path, format="PNG")
                    img = fitz.Pixmap(frame_path)
                    rect = fitz.Rect(0, 0, img.width, img.height)
                    page = doc.new_page(width=rect.width, height=rect.height)
                    page.insert_image(rect, pixmap=img)
                    os.remove(frame_path)
                doc.save(pdf_path)
                doc.close()
                clip.close()
                self.file_handler.post_process(movie_path_set, pdf_path, self.delete)
        # Convert Documents to PDF
        for doc_path_set in file_paths[Category.DOCUMENT]:
            if doc_path_set[2] == "pdf":
                # If the document is already a pdf, skip it
                continue
            if doc_path_set[2] == "srt":
                # Convert srt to pdf
                pdf_path = os.path.abspath(
                    os.path.join(self.output, f"{doc_path_set[1]}.{format}")
                )
                with open(self.file_handler.join_back(doc_path_set), "r") as srt_file:
                    srt_content = srt_file.read()
                # Insert the SRT content into the PDF
                doc = fitz.open()
                page = doc.new_page()
                page.insert_text((50, 50), srt_content, fontsize=12)
                doc.save(pdf_path)
                doc.close()
                self.file_handler.post_process(doc_path_set, pdf_path, self.delete)
            elif doc_path_set[2] == "docx":
                pdf_path = os.path.abspath(
                    os.path.join(self.output, f"{doc_path_set[1]}.{format}")
                )
                docx_doc = open(self.file_handler.join_back(doc_path_set), "rb")
                # Convert docx to HTML as intermediary
                document = mammoth.convert_to_html(docx_doc)
                docx_doc.close()
                # Convert html to PDF, save that
                HTML(string=document.value.encode("utf-8")).write_pdf(pdf_path)
                self.file_handler.post_process(doc_path_set, pdf_path, self.delete)

    def to_frames(self, file_paths: dict, format: str) -> None:
        # Converting to image frame sets
        # This works for images and movies only
        format = "jpeg" if format == "jpg" else format
        gifs_to_frames(self.output, file_paths, self.file_handler)
        for image_path_set in file_paths[Category.IMAGE]:
            if image_path_set[2] == format:
                continue
            if image_path_set[2] == "gif":
                # gifs_to_frames did that out of loop already, just logging here
                self.file_handler.post_process(image_path_set, self.output, self.delete)
            else:
                if not os.path.exists(os.path.join(self.output, image_path_set[1])):
                    self.output = self.input
                img_path = os.path.abspath(
                    os.path.join(self.output, f"{image_path_set[1]}.{format}")
                )
                with Image.open(self.file_handler.join_back(image_path_set)) as img:
                    img.convert("RGB").save(img_path, format=format)
                self.file_handler.post_process(image_path_set, img_path, self.delete)
        # Convert documents to image frames
        for doc_path_set in file_paths[Category.DOCUMENT]:
            if doc_path_set[2] == format:
                continue
            if not os.path.exists(os.path.join(self.output, doc_path_set[1])):
                try:
                    os.makedirs(
                        os.path.join(self.output, doc_path_set[1]), exist_ok=True
                    )
                except OSError as e:
                    self.event_logger.info(
                        f"[!] {lang.get_translation('error', self.locale)}: {e} - {lang.get_translation('set_out_dir', self.locale)} {self.input}"
                    )
                    self.output = self.input
            if doc_path_set[2] in ["docx", "pptx"]:
                # Read all images from docx, write to os.path.join(self.output, doc_path_set[1])
                office_to_frames(
                    doc_path_set=doc_path_set,
                    format=format,
                    output=self.output,
                    delete=self.delete,
                    file_handler=self.file_handler,
                    event_logger=self.event_logger,
                )
            if doc_path_set[2] == "pdf":
                # Per page, convert pdf to image
                pdf_path = self.file_handler.join_back(doc_path_set)
                pdf = PyPDF2.PdfReader(pdf_path)
                img_path = os.path.abspath(
                    os.path.join(
                        os.path.join(self.output, doc_path_set[1]),
                        f"{doc_path_set[1]}-%{len(str(len(pdf.pages)))}d.{format}",
                    )
                )

                try:
                    if not os.path.exists(os.path.dirname(img_path)):
                        os.makedirs(os.path.dirname(img_path), exist_ok=True)
                    self.output = os.path.dirname(img_path)
                except OSError as e:
                    self.event_logger.info(
                        f"[!] {lang.get_translation('error', self.locale)}: {e} - {lang.get_translation('set_out_dir', self.locale)} {self.input}"
                    )
                    self.output = self.input

                pdf_document = fitz.open(pdf_path)

                for page_num in range(len(pdf_document)):
                    pix = pdf_document[page_num].get_pixmap()
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    img_file = img_path % (page_num + 1)
                    img.save(img_file, format.upper())

                pdf_document.close()
                self.file_handler.post_process(doc_path_set, img_path, self.delete)

        # Audio cant be image-framed, movies certrainly can
        for movie_path_set in file_paths[Category.MOVIE]:
            if movie_path_set[2] not in self._supported_formats[Category.MOVIE]:
                self.event_logger.info(
                    f"[!] {lang.get_translation('movie_format_unsupported', self.locale)} {movie_path_set[2]} - {lang.get_translation('skipping', self.locale)}"
                )
                continue
            if self.file_handler.has_visuals(movie_path_set):
                video = VideoFileClip(
                    self.file_handler.join_back(movie_path_set),
                    audio=False,
                    fps_source="tbr",
                )
                try:
                    if not os.path.exists(os.path.join(self.output, movie_path_set[1])):
                        os.makedirs(
                            os.path.join(self.output, movie_path_set[1]), exist_ok=True
                        )
                except OSError as e:
                    self.event_logger.info(
                        f"[!] {lang.get_translation('error', self.locale)}: {e} - {lang.get_translation('set_out_dir', self.locale)} {self.input}"
                    )
                    self.output = self.input
                img_path = os.path.abspath(
                    os.path.join(
                        os.path.join(self.output, movie_path_set[1]),
                        f"{movie_path_set[1]}-%{len(str(int(video.duration * video.fps)))}d.{format}",
                    )
                )

                video.write_images_sequence(
                    img_path, fps=video.fps, logger=self.prog_logger
                )
                video.close()
                self.file_handler.post_process(movie_path_set, img_path, self.delete)
            else:
                self.event_logger.info(
                    f'[!] {lang.get_translation("skipping", self.locale)} "{self.file_handler.join_back(movie_path_set)}" - {lang.get_translation("audio_only_video", self.locale)}'
                )
                self.file_handler.post_process(movie_path_set, img_path, self.delete)

    def to_gif(self, file_paths: dict, format: str) -> None:
        # All images in the input directory are merged into one gif
        if len(file_paths[Category.IMAGE]) > 0:
            images = []
            for image_path_set in file_paths[Category.IMAGE]:
                if image_path_set[2] == format:
                    continue
                with Image.open(self.file_handler.join_back(image_path_set)) as image:
                    images.append(image.convert("RGB"))
            images[0].save(
                os.path.join(self.output, f"merged.{format}"),
                save_all=True,
                append_images=images[1:],
            )

        # Movies are converted to gifs as well, retaining 1/3 of the frames
        for movie_path_set in file_paths[Category.MOVIE]:
            if self.file_handler.has_visuals(movie_path_set):
                video = VideoFileClip(
                    self.file_handler.join_back(movie_path_set),
                    audio=False,
                    fps_source="tbr",
                )
                gif_path = os.path.join(self.output, f"{movie_path_set[1]}.{format}")
                video.write_gif(gif_path, fps=video.fps // 3, logger=self.prog_logger)
                video.close()
                self.file_handler.post_process(movie_path_set, gif_path, self.delete)
            else:
                self.event_logger.info(
                    f'[!] {lang.get_translation("skipping", self.locale)} "{self.file_handler.join_back(movie_path_set)}" - {lang.get_translation("audio_only_video", self.locale)}'
                )
        # Documents may be convertable to gifs, e.g. pdfs
        for doc_path_set in file_paths[Category.DOCUMENT]:
            if doc_path_set[2] == "pdf":
                pdf_path = self.file_handler.join_back(doc_path_set)
                gif_path = os.path.abspath(
                    os.path.join(self.output, f"{doc_path_set[1]}.{format}")
                )
                doc = fitz.open(pdf_path)
                images = []
                for page_num in range(len(doc)):
                    pix = doc.load_page(page_num).get_pixmap()
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    images.append(img.convert("RGB"))
                if images:
                    images[0].save(
                        gif_path,
                        save_all=True,
                        append_images=images[1:],
                        duration=(len(images) * 1000 // len(doc))
                        // (12 if self.framerate is None else self.framerate),
                        loop=0,
                    )
                doc.close()
                self.file_handler.post_process(doc_path_set, gif_path, self.delete)
            elif doc_path_set[2] in ["docx", "pptx"]:
                input_path = self.file_handler.join_back(doc_path_set)
                gif_path = os.path.abspath(
                    os.path.join(self.output, f"{doc_path_set[1]}.{format}")
                )
                images = []
                if doc_path_set[2] == "docx":
                    doc = docx.Document(input_path)
                    for rel in doc.part.rels.values():
                        if "image" in rel.reltype:
                            img = Image.open(rel.target_part.blob)
                            images.append(img.convert("RGB"))
                    frame_count = len(doc.paragraphs) or 1
                else:
                    prs = pptx.Presentation(input_path)
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if shape.shape_type == 13:  # Picture
                                image = shape.image
                                img_bytes = image.blob
                                img = Image.open(img_bytes)
                                images.append(img.convert("RGB"))
                    frame_count = len(prs.slides) or 1
                if images:
                    images[0].save(
                        gif_path,
                        save_all=True,
                        append_images=images[1:],
                        duration=(len(images) * 1000 // frame_count)
                        // (12 if self.framerate is None else self.framerate),
                        loop=0,
                    )
                self.file_handler.post_process(doc_path_set, gif_path, self.delete)

    def to_bmp(self, file_paths: dict, format: str) -> None:
        for movie_path_set in file_paths[Category.MOVIE]:
            # Movies are converted to bmps frame by frame
            if self.file_handler.has_visuals(movie_path_set):
                video = VideoFileClip(
                    self.file_handler.join_back(movie_path_set),
                    audio=False,
                    fps_source="tbr",
                )
                bmp_path = os.path.join(self.output, f"{movie_path_set[1]}.{format}")
                # Split video into individual bmp frame images at original framerate
                for _, frame in enumerate(
                    video.iter_frames(fps=video.fps, dtype="uint8")
                ):
                    frame.save(
                        f"{bmp_path}-%{len(str(int(video.duration * video.fps)))}d.{format}",
                        format=format,
                    )
                self.file_handler.post_process(movie_path_set, bmp_path, self.delete)
            else:
                self.event_logger.info(
                    f'[!] {lang.get_translation("skipping", self.locale)} "{self.file_handler.join_back(movie_path_set)}" - {lang.get_translation("audio_only_video", self.locale)}'
                )
        for image_path_set in file_paths[Category.IMAGE]:
            # Pngs and gifs are converted to bmps as well
            if image_path_set[2] == format:
                continue
            if image_path_set[2] in ["png", "jpeg", "jpg", "tiff", "tga", "eps"]:
                bmp_path = os.path.join(self.output, f"{image_path_set[1]}.{format}")
                with Image.open(self.file_handler.join_back(image_path_set)) as img:
                    img.convert("RGB").save(bmp_path, format=format)
                self.file_handler.post_process(image_path_set, bmp_path, self.delete)
            elif image_path_set[2] == "gif":
                clip = VideoFileClip(self.file_handler.join_back(image_path_set))
                for _, frame in enumerate(
                    clip.iter_frames(fps=clip.fps, dtype="uint8")
                ):
                    frame_path = os.path.join(
                        self.output,
                        f"{image_path_set[1]}-%{len(str(int(clip.duration * clip.fps)))}d.{format}",
                    )
                    Image.fromarray(frame).convert("RGB").save(
                        frame_path, format=format
                    )
                self.file_handler.post_process(image_path_set, frame_path, self.delete)
            else:
                # Handle unsupported file types here
                self.event_logger.info(
                    f'[!] {lang.get_translation("skipping", self.locale)} "{self.file_handler.join_back(image_path_set)}" - {lang.get_translation("unsupported_format", self.locale)}'
                )
        # Documents can sometimes be converted to bmp, e.g. contents of docx, pdf
        for doc_path_set in file_paths[Category.DOCUMENT]:
            if doc_path_set[2] == "docx":
                office_to_frames(
                    doc_path_set=doc_path_set,
                    format=format,
                    output=self.output,
                    delete=self.delete,
                    file_handler=self.file_handler,
                    event_logger=self.event_logger,
                )
            if doc_path_set[2] == "pdf":
                pdf_path = self.file_handler.join_back(doc_path_set)
                bmp_path = os.path.abspath(
                    os.path.join(self.output, f"{doc_path_set[1]}.{format}")
                )
                doc = fitz.open(pdf_path)
                if not os.path.exists(os.path.join(self.output, doc_path_set[1])):
                    try:
                        os.makedirs(
                            os.path.join(self.output, doc_path_set[1]), exist_ok=True
                        )
                    except OSError as e:
                        self.event_logger.info(
                            f"[!] {lang.get_translation('error', self.locale)}: {e} - {lang.get_translation('set_out_dir', self.locale)} {self.input}"
                        )
                        self.output = self.input
                for i, page_num in enumerate(range(len(doc))):
                    pix = doc.load_page(page_num).get_pixmap()
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    img = img.convert("RGB")
                    # Save each page as a separate BMP file
                    img.save(
                        os.path.join(
                            self.output,
                            doc_path_set[1],
                            f"{doc_path_set[1]}-{i:0{len(str(len(doc)))}}.{format}",
                        ),
                        format=format,
                    )
                doc.close()
                self.file_handler.post_process(doc_path_set, bmp_path, self.delete)

    def to_webp(self, file_paths: dict, format: str) -> None:
        # Convert frames in webp format
        # Movies are converted to webps, frame by frame
        for movie_path_set in file_paths[Category.MOVIE]:
            if self.file_handler.has_visuals(movie_path_set):
                video = VideoFileClip(
                    self.file_handler.join_back(movie_path_set),
                    audio=False,
                    fps_source="tbr",
                )
                if not os.path.exists(os.path.join(self.output, movie_path_set[1])):
                    try:
                        os.makedirs(os.path.join(self.output, movie_path_set[1]))
                    except OSError as e:
                        self.event_logger.info(
                            f"[!] {lang.get_translation('error', self.locale)}: {e} - {lang.get_translation('set_out_dir', self.locale)} {self.input}"
                        )
                        self.output = self.input
                img_path = os.path.abspath(
                    os.path.join(
                        os.path.join(self.output, movie_path_set[1]),
                        f"{movie_path_set[1]}-%{len(str(int(video.duration * video.fps)))}d.{format}",
                    )
                )
                video.write_images_sequence(
                    img_path, fps=video.fps, logger=self.prog_logger
                )
                video.close()
                self.file_handler.post_process(movie_path_set, img_path, self.delete)
            else:
                self.event_logger.info(
                    f'[!] {lang.get_translation("skipping", self.locale)} "{self.file_handler.join_back(movie_path_set)}" - {lang.get_translation("audio_only_video", self.locale)}'
                )

        # pngs and gifs are converted to webps as well
        for image_path_set in file_paths[Category.IMAGE]:
            if image_path_set[2] == format:
                continue
            if image_path_set[2] in ["png", "jpeg", "jpg", "tiff", "tga", "eps"]:
                webp_path = os.path.join(self.output, f"{image_path_set[1]}.{format}")
                with Image.open(self.file_handler.join_back(image_path_set)) as img:
                    img.convert("RGB").save(webp_path, format=format)
                self.file_handler.post_process(image_path_set, webp_path, self.delete)
            elif image_path_set[2] == "gif":
                clip = VideoFileClip(self.file_handler.join_back(image_path_set))
                for _, frame in enumerate(
                    clip.iter_frames(fps=clip.fps, dtype="uint8")
                ):
                    frame_path = os.path.join(
                        self.output,
                        f"{image_path_set[1]}-%{len(str(int(clip.duration * clip.fps)))}d.{format}",
                    )
                    Image.fromarray(frame).convert("RGB").save(
                        frame_path, format=format
                    )
                self.file_handler.post_process(image_path_set, frame_path, self.delete)
            else:
                # Handle unsupported file types here
                self.event_logger.info(
                    f'[!] {lang.get_translation("skipping", self.locale)} "{self.file_handler.join_back(image_path_set)}" - {lang.get_translation("unsupported_format", self.locale)}'
                )

        # Documents may be convertable to bmps, e.g. pdfs
        for doc_path_set in file_paths[Category.DOCUMENT]:
            if doc_path_set[2] == "docx":
                office_to_frames(
                    doc_path_set=doc_path_set,
                    format=format,
                    output=self.output,
                    delete=self.delete,
                    file_handler=self.file_handler,
                    event_logger=self.event_logger,
                )
            if doc_path_set[2] == "pdf":
                pdf_path = self.file_handler.join_back(doc_path_set)
                bmp_path = os.path.abspath(
                    os.path.join(self.output, f"{doc_path_set[1]}.{format}")
                )
                doc = fitz.open(pdf_path)
                if not os.path.exists(os.path.join(self.output, doc_path_set[1])):
                    try:
                        os.makedirs(
                            os.path.join(self.output, doc_path_set[1]), exist_ok=True
                        )
                    except OSError as e:
                        self.event_logger.info(
                            f"[!] {lang.get_translation('error', self.locale)}: {e} - {lang.get_translation('set_out_dir', self.locale)} {self.input}"
                        )
                        self.output = self.input
                for i, page_num in enumerate(range(len(doc))):
                    pix = doc.load_page(page_num).get_pixmap()
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    img = img.convert("RGB")
                    # Save each page as a separate BMP file
                    img.save(
                        os.path.join(
                            self.output,
                            doc_path_set[1],
                            f"{doc_path_set[1]}-{i:0{len(str(len(doc)))}}.{format}",
                        ),
                        format=format,
                    )
                doc.close()
                self.file_handler.post_process(doc_path_set, bmp_path, self.delete)

    def concat(self, file_paths: dict, format: str) -> None:
        # Concatenate files of same type (img/movie/audio) back to back
        # Concatenate audio files
        if file_paths[Category.AUDIO] and (
            format is None or format in self._supported_formats[Category.AUDIO]
        ):
            concat_audio = concatenate_audioclips(
                [
                    AudioFileClip(self.file_handler.join_back(audio_path_set))
                    for audio_path_set in file_paths[Category.AUDIO]
                ]
            )
            format = "mp3" if format is None else format
            audio_out_path = os.path.join(self.output, f"concatenated_audio.{format}")
            concat_audio.write_audiofile(
                audio_out_path,
                codec=self._supported_formats[Category.AUDIO][format],
                bitrate=self._audio_bitrate(format, self.quality)
                if self.quality is not None
                else getattr(concat_audio, "bitrate", "192k"),
                logger=self.prog_logger,
            )
            concat_audio.close()
        # Concatenate movie files
        if file_paths[Category.MOVIE] and (
            format is None or format in self._supported_formats[Category.MOVIE]
        ):
            concat_vid = concatenate_videoclips(
                [
                    VideoFileClip(
                        self.file_handler.join_back(movie_path_set),
                        audio=True,
                        fps_source="tbr",
                    )
                    for movie_path_set in file_paths[Category.MOVIE]
                ],
                method="compose",
            )
            format = "mp4" if format is None else format
            video_out_path = os.path.join(self.output, f"concatenated_video.{format}")
            concat_vid.write_videofile(
                video_out_path,
                fps=concat_vid.fps if self.framerate is None else concat_vid.fps,
                codec=self._supported_formats[Category.MOVIE][format],
                logger=self.prog_logger,
            )
            concat_vid.close()
        # Concatenate image files (make a gif out of them)
        if file_paths[Category.IMAGE] and (
            format is None or format in self._supported_formats[Category.IMAGE]
        ):
            gif_out_path = os.path.join(self.output, "concatenated_image.gif")
            concatenated_image = clips_array(
                [
                    [
                        ImageClip(
                            self.file_handler.join_back(image_path_set)
                        ).with_duration(
                            1 / 12 if self.framerate is None else 1 / self.framerate
                        )
                    ]
                    for image_path_set in file_paths[Category.IMAGE]
                ]
            )
            concatenated_image.write_gif(
                gif_out_path, fps=self.framerate, logger=self.prog_logger
            )
        # Concatenate document files (keep respective document format)
        if file_paths[Category.DOCUMENT] and (
            format is None or format in self._supported_formats[Category.DOCUMENT]
        ):
            pdf_out_path = os.path.join(self.output, "concatenated_pdfs.pdf")
            pdfs = sorted(
                [
                    doc_path_set if doc_path_set[2] == "pdf" else None
                    for doc_path_set in file_paths[Category.DOCUMENT]
                ],
                key=lambda x: x[1],
            )
            srt_out_path = os.path.join(self.output, "concatenated_subtitles.srt")
            srts = sorted(
                [
                    doc_path_set if doc_path_set[2] == "srt" else None
                    for doc_path_set in file_paths[Category.DOCUMENT]
                ],
                key=lambda x: x[1],
            )
            if len(pdfs) > 0:
                # Produce a single pdf file
                doc = fitz.open()
                for doc_path_set in pdfs:
                    pdf_path = self.file_handler.join_back(doc_path_set)
                    pdf_document = fitz.open(pdf_path)
                    doc.insert_pdf(pdf_document)
                    pdf_document.close()
                doc.save(pdf_out_path)
                doc.close()
            if len(srts) > 0:
                for doc_path_set in srts:
                    # Produce a single srt file
                    srt_path = self.file_handler.join_back(doc_path_set)
                    with open(srt_path, "r") as srt_file:
                        srt_content = srt_file.read()
                    # Insert the SRT content into the PDF
                    with open(srt_out_path, "a") as srt_file:
                        srt_file.write(srt_content)
                        srt_file.write("\n")
        for category in file_paths.keys():
            # Iterate over each input category and post-process respective files
            for i, file_path in enumerate(file_paths[category]):
                self.file_handler.post_process(
                    file_path, self.output, self.delete, show_status=(i == 0)
                )
        self.event_logger.info(
            f"[+] {lang.get_translation('concat_success', self.locale)}"
        )

    def merge(self, file_paths: dict, across: bool = False) -> None:
        # For movie files and equally named audio file, merge them together under same name
        # (movie with audio with '_merged' addition to name)
        # If only a video file is provided, look for a matching audio file in the same directory
        found_audio = False
        audio_exts = list(self._supported_formats[Category.AUDIO].keys())

        for movie_path_set in file_paths[Category.MOVIE]:
            # Try to find a corresponding audio file in the input set
            # (e.g. "-1 path1 -2 path2 -n pathn")
            if across:
                # Allow matching audio from any input directory
                audio_fit = next(
                    (
                        audio_set
                        for audio_set in file_paths[Category.AUDIO]
                        if audio_set[1] == movie_path_set[1]
                    ),
                    None,
                )
            else:
                # Only match audio from the same directory as the video
                audio_fit = next(
                    (
                        audio_set
                        for audio_set in file_paths[Category.AUDIO]
                        if audio_set[1] == movie_path_set[1]
                        and audio_set[0] == movie_path_set[0]
                    ),
                    None,
                )

            # If not found, look for a matching audio file in the video's directory
            if audio_fit is None:
                video_dir = movie_path_set[0]
                video_basename = movie_path_set[1]
                for ext in audio_exts:
                    candidate = os.path.join(video_dir, f"{video_basename}.{ext}")
                    if os.path.isfile(candidate):
                        audio_fit = (video_dir, video_basename, ext)
                        break

            if audio_fit is not None:
                found_audio = True
                # Merge movie and audio file
                audio = None
                video = None
                try:
                    video = VideoFileClip(self.file_handler.join_back(movie_path_set))
                    audio = AudioFileClip(self.file_handler.join_back(audio_fit))
                    video = video.set_audio(audio)
                    merged_out_path = os.path.join(
                        self.output, f"{movie_path_set[1]}_merged.{movie_path_set[2]}"
                    )
                    video.write_videofile(
                        merged_out_path,
                        fps=video.fps if self.framerate is None else self.framerate,
                        codec=self._supported_formats[Category.MOVIE][
                            movie_path_set[2]
                        ],
                        logger=self.prog_logger,
                    )
                finally:
                    if audio is not None:
                        audio.close()
                    if video is not None:
                        video.close()
                self.file_handler.post_process(
                    movie_path_set, merged_out_path, self.delete
                )
                # Only delete the audio file if it was in the input set, not if just found in dir
                if audio_fit in file_paths[Category.AUDIO]:
                    self.file_handler.post_process(
                        audio_fit, merged_out_path, self.delete, show_status=False
                    )

        if not found_audio:
            self.event_logger.warning(
                f"[!] {lang.get_translation('no_audio_movie_match', self.locale)}"
            )

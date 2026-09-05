import os
import sys
import docx
import pptx
import pymupdf as pypdf
import utils.language_support as lang

from PIL import Image
from tqdm import tqdm
from io import BytesIO
from moviepy import VideoFileClip
from utils.category import Category
from concurrent.futures import ThreadPoolExecutor, as_completed
from core.utils.file_handler import safe_output_dir as _safe_out_dir


sys.path.insert(0, os.path.abspath(os.path.join(os.path.dirname(__file__), "..")))

try:
    import pillow_heif

    pillow_heif.register_heif_opener()
except ImportError:
    pillow_heif = None

_SAVE_FORMAT = {
    "jpg": "JPEG",
    "jpeg": "JPEG",
    "heic": "HEIF",
    "heif": "HEIF",
    "avif": "AVIF",
}


def _save_format(format: str) -> str:
    if format in _SAVE_FORMAT:
        return _SAVE_FORMAT[format]
    return format.upper()


def _write_office_frame(img_bytes: bytes, out_path: str, format: str) -> None:
    target = _save_format(format)
    try:
        with Image.open(BytesIO(img_bytes)) as img:
            if img.format != target:
                img.convert("RGB").save(out_path, format=target)
                return
    except Exception:
        pass
    with open(out_path, "wb") as f:
        f.write(img_bytes)


def office_to_frames(
    doc_path_set: tuple,
    format: str,
    output: str,
    delete: bool,
    file_handler,
    event_logger,
) -> str | None:
    # Extract imgs from office (docx, pptx), save as frames in specified format
    full_path = file_handler.join_back(doc_path_set)
    out_dir = _safe_out_dir(os.path.join(output, doc_path_set[1]))
    out_path = None
    os.makedirs(out_dir, exist_ok=True)
    try:
        if doc_path_set[2] == "docx":
            doc = docx.Document(full_path)
            for i, rel in tqdm(
                enumerate(doc.part.rels.values()), desc=f"{doc_path_set[1]}"
            ):
                if "image" in rel.reltype:
                    img_bytes = rel.target_part.blob
                    out_path = os.path.join(out_dir, f"{doc_path_set[1]}-{i}.{format}")
                    _write_office_frame(img_bytes, out_path, format)
            if out_path is not None:
                file_handler.post_process(doc_path_set, out_path, delete)
        elif doc_path_set[2] == "pptx":
            prs = pptx.Presentation(full_path)
            img_count = 0
            for i, slide in tqdm(enumerate(prs.slides), desc=f"{doc_path_set[1]}"):
                for shape in slide.shapes:
                    if shape.shape_type == 13:  # Picture type
                        image = shape.image
                        img_bytes = image.blob
                        out_path = os.path.join(
                            out_dir, f"{doc_path_set[1]}-{img_count}.{format}"
                        )
                        _write_office_frame(img_bytes, out_path, format)
                        img_count += 1
            if out_path is not None:
                file_handler.post_process(doc_path_set, out_path, delete)
    except Exception as e:
        event_logger.error(e)
    return out_dir


def _max_workers() -> int:
    try:
        _cpu = os.cpu_count()
        _cpu = _cpu if isinstance(_cpu, int) else 2
        return max(1, min(int(os.environ.get("Any2Any_MAX_WORKERS", "1")), _cpu - 1))
    except (ValueError, TypeError):
        return 1


def gif_to_frames(output: str, file_paths: dict, file_handler) -> dict:
    # Convert GIFs to frames, place those in a folder
    gifs = [
        image_path
        for image_path in file_paths[Category.IMAGE]
        if image_path[2] == "gif"
    ]

    if not gifs:
        return {}

    created: dict = {}

    def _extract_gif_frames(image_path_set: tuple):
        base_dir = os.path.abspath(os.path.join(output, image_path_set[1]))
        out_dir = _safe_out_dir(base_dir)
        os.makedirs(out_dir, exist_ok=True)
        created[image_path_set[1]] = out_dir
        clip = VideoFileClip(file_handler.join_back(image_path_set), audio=False)
        try:
            total_frames = int(clip.duration * clip.fps)
            num_digits = len(str(total_frames))
            for i, frame in enumerate(clip.iter_frames(fps=clip.fps, dtype="uint8")):
                Image.fromarray(frame).save(
                    os.path.join(out_dir, f"{image_path_set[1]}-{i:0{num_digits}d}.png")
                )
        finally:
            clip.close()

    max_workers = _max_workers()
    if len(gifs) == 1 or max_workers == 1:
        for g in gifs:
            _extract_gif_frames(g)
    else:
        with ThreadPoolExecutor(max_workers=max_workers) as ex:
            list(ex.map(_extract_gif_frames, gifs))
    return created


def _save_clip_frames(
    clip: VideoFileClip, out_dir: str, base_name: str, fmt: str
) -> str | None:
    # Return path of last frame written, or None for a clip without frames
    os.makedirs(out_dir, exist_ok=True)
    num_digits = len(str(int(clip.duration * clip.fps)))
    last_frame = None
    for i, frame in enumerate(clip.iter_frames(fps=clip.fps, dtype="uint8")):
        last_frame = os.path.join(out_dir, f"{base_name}-{i:0{num_digits}d}.{fmt}")
        Image.fromarray(frame).convert("RGB").save(
            last_frame, format=_save_format(fmt)
        )
    return last_frame


class ImageConverter:
    def __init__(
        self, file_handler, prog_logger, event_logger, locale: str = "English"
    ):
        self.file_handler = file_handler
        self.prog_logger = prog_logger
        self.event_logger = event_logger
        self.locale = locale

    def _per_file_out_dir(self, path_set: tuple, input: str, output: str) -> str:
        # When in_dir doubles as out_dir: place outputs next to source
        if input == output:
            return path_set[0]
        return output

    def to_frames(
        self,
        input: str,
        output: str,
        file_paths: dict,
        supported_formats: dict,  # self._supported_formats
        framerate: int,
        format: str,
        delete: bool,
    ) -> None:
        # Converting to image frame sets
        # This works for images and movies only
        format = "jpeg" if format == "jpg" else format

        # Handle GIF files first
        gif_to_frames(output, file_paths, self.file_handler)

        # Process regular image files
        for image_path_set in file_paths[Category.IMAGE]:
            try:
                if image_path_set[2] == format:
                    continue

                if image_path_set[2] == "gif":
                    # gif_to_frames already processed these
                    self.file_handler.post_process(image_path_set, output, delete)
                    continue

                # Ensure output dir exists (single images are placed directly in output)
                os.makedirs(output, exist_ok=True)

                out_dir = self._per_file_out_dir(image_path_set, input, output)
                img_path = self.file_handler._resolve_output_file_conflict(
                    os.path.abspath(
                        os.path.join(out_dir, f"{image_path_set[1]}.{format}")
                    )
                )

                # Convert and save the image
                with Image.open(self.file_handler.join_back(image_path_set)) as img:
                    # Convert to RGB if needed (important for WebP to PNG)
                    if img.mode != "RGB":
                        img = img.convert("RGB")
                    img.save(img_path, format=_save_format(format))

                # Verify the file was created before deleting source
                if os.path.exists(img_path):
                    self.file_handler.post_process(image_path_set, img_path, delete)
                else:
                    self.event_logger.error(
                        f"[!] {lang.get_translation('conversion_failed', self.locale)}: "
                        f"{self.file_handler.join_back(image_path_set)}"
                    )

            except Exception as e:
                self.event_logger.error(
                    f"[!] {lang.get_translation('error', self.locale)} "
                    f"converting {self.file_handler.join_back(image_path_set)}: {str(e)}"
                )
                # Don't delete source file if conversion failed
                delete = False
                continue
        # Convert documents to image frames
        for doc_path_set in file_paths[Category.DOCUMENT]:
            if doc_path_set[2] == format:
                continue
            _base_doc_dir = (
                os.path.join(doc_path_set[0], doc_path_set[1])
                if input == output
                else os.path.join(output, doc_path_set[1])
            )
            doc_out_dir = _safe_out_dir(_base_doc_dir)
            if not os.path.exists(doc_out_dir):
                try:
                    os.makedirs(doc_out_dir, exist_ok=True)
                except OSError as e:
                    self.event_logger.info(
                        f"[!] {lang.get_translation('error', self.locale)}: {e} - {lang.get_translation('set_out_dir', self.locale)} {input}"
                    )
                    output = input
                    doc_out_dir = _safe_out_dir(
                        os.path.join(output, doc_path_set[1])
                    )
                    os.makedirs(doc_out_dir, exist_ok=True)
            if doc_path_set[2] in ["docx", "pptx"]:
                # Read all images from docx, write to os.path.join(self.output, doc_path_set[1])
                office_to_frames(
                    doc_path_set=doc_path_set,
                    format=format,
                    output=output,
                    delete=delete,
                    file_handler=self.file_handler,
                    event_logger=self.event_logger,
                )
            if doc_path_set[2] == "pdf":
                pdf_path = self.file_handler.join_back(doc_path_set)
                pdf_document = pypdf.open(pdf_path)
                total_pages = len(pdf_document)
                img_path_pattern = os.path.abspath(
                    os.path.join(
                        doc_out_dir,
                        f"{doc_path_set[1]}-%0{len(str(total_pages))}d.{format}",
                    )
                )
                img_file = None
                for page_num in range(total_pages):
                    pix = pdf_document[page_num].get_pixmap()
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    img_file = img_path_pattern % (page_num + 1)
                    img.save(img_file, _save_format(format))
                pdf_document.close()
                if img_file is not None:
                    self.file_handler.post_process(doc_path_set, img_file, delete)

        # Audio cant be image-framed, movies certrainly can
        def _movie_to_frames(movie_path_set: tuple):
            if movie_path_set[2] not in supported_formats[Category.MOVIE]:
                self.event_logger.info(
                    f"[!] {lang.get_translation('movie_format_unsupported', self.locale)} {movie_path_set[2]} - {lang.get_translation('skipping', self.locale)}"
                )
                return None
            if not self.file_handler.has_visuals(movie_path_set):
                self.event_logger.info(
                    f'[!] {lang.get_translation("skipping", self.locale)} "{self.file_handler.join_back(movie_path_set)}" - {lang.get_translation("audio_only_video", self.locale)}'
                )
                return None
            video = VideoFileClip(
                self.file_handler.join_back(movie_path_set),
                audio=False,
                fps_source="tbr",
            )
            try:
                movie_out_dir = _safe_out_dir(
                    os.path.join(movie_path_set[0], movie_path_set[1])
                    if input == output
                    else os.path.join(output, movie_path_set[1])
                )
                try:
                    os.makedirs(movie_out_dir, exist_ok=True)
                except OSError as e:
                    self.event_logger.info(
                        f"[!] {lang.get_translation('error', self.locale)}: {e} - {lang.get_translation('set_out_dir', self.locale)} {input}"
                    )
                    return None
                img_path = os.path.abspath(
                    os.path.join(
                        movie_out_dir,
                        f"{movie_path_set[1]}-%0{len(str(int(video.duration * video.fps)))}d.{format}",
                    )
                )
                frame_names = video.write_images_sequence(
                    img_path, fps=video.fps, logger=self.prog_logger
                )
                return (
                    (movie_path_set, frame_names[-1]) if frame_names else None
                )
            finally:
                video.close()

        max_workers = _max_workers()
        movie_items = list(file_paths[Category.MOVIE])
        if len(movie_items) <= 1 or max_workers == 1:
            for m in movie_items:
                res = _movie_to_frames(m)
                if res is not None:
                    self.file_handler.post_process(res[0], res[1], delete)
        else:
            with ThreadPoolExecutor(max_workers=max_workers) as ex:
                futures = [ex.submit(_movie_to_frames, m) for m in movie_items]
                for fut in as_completed(futures):
                    res = fut.result()
                    if res is not None:
                        self.file_handler.post_process(res[0], res[1], delete)

    def to_bmp(
        self,
        input: str,
        output: str,
        file_paths: dict,
        supported_formats: dict,  # unused, aligns signature with to_frames, so keep it
        framerate: int,
        format: str,
        delete: bool,
    ) -> None:
        for movie_path_set in file_paths[Category.MOVIE]:
            # Movies are converted to bmps frame by frame
            if self.file_handler.has_visuals(movie_path_set):
                video = VideoFileClip(
                    self.file_handler.join_back(movie_path_set),
                    audio=False,
                    fps_source="tbr",
                )
                try:
                    movie_bmp_out_dir = _safe_out_dir(
                        os.path.join(movie_path_set[0], movie_path_set[1])
                        if input == output
                        else os.path.join(output, movie_path_set[1])
                    )
                    # Split video into individual bmp frame images at original framerate
                    bmp_frame_path = _save_clip_frames(
                        video, movie_bmp_out_dir, movie_path_set[1], format
                    )
                finally:
                    video.close()
                if bmp_frame_path is not None:
                    self.file_handler.post_process(movie_path_set, bmp_frame_path, delete)
            else:
                self.event_logger.info(
                    f'[!] {lang.get_translation("skipping", self.locale)} "{self.file_handler.join_back(movie_path_set)}" - {lang.get_translation("audio_only_video", self.locale)}'
                )

        convertible_images = [
            ips
            for ips in file_paths[Category.IMAGE]
            if ips[2] != format
            and ips[2]
            in [
                "png",
                "jpeg",
                "jpg",
                "tiff",
                "tga",
                "eps",
                "heic",
                "heif",
                "avif",
            ]
        ]

        def _convert_image_to_bmp(image_path_set: tuple):
            out_dir = self._per_file_out_dir(image_path_set, input, output)
            bmp_path = self.file_handler._resolve_output_file_conflict(
                os.path.abspath(os.path.join(out_dir, f"{image_path_set[1]}.{format}"))
            )
            with Image.open(self.file_handler.join_back(image_path_set)) as img:
                img.convert("RGB").save(bmp_path, format=format)
            return (image_path_set, bmp_path)

        max_workers = _max_workers()
        if len(convertible_images) <= 1 or max_workers == 1:
            for image_path_set in convertible_images:
                src, bmp_path = _convert_image_to_bmp(image_path_set)
                self.file_handler.post_process(src, bmp_path, delete)
        else:
            with ThreadPoolExecutor(max_workers=max_workers) as ex:
                futures = [
                    ex.submit(_convert_image_to_bmp, ips) for ips in convertible_images
                ]
                for fut in as_completed(futures):
                    src, bmp_path = fut.result()
                    self.file_handler.post_process(src, bmp_path, delete)

        for image_path_set in file_paths[Category.IMAGE]:
            # Pngs and gifs are converted to bmps as well
            if image_path_set[2] == format:
                continue
            if image_path_set[2] in [
                "png",
                "jpeg",
                "jpg",
                "tiff",
                "tga",
                "eps",
                "heic",
                "heif",
                "avif",
            ]:
                continue
            elif image_path_set[2] == "gif":
                clip = VideoFileClip(self.file_handler.join_back(image_path_set))
                try:
                    gif_frame_out_dir = _safe_out_dir(
                        os.path.abspath(os.path.join(output, image_path_set[1]))
                    )
                    frame_path = _save_clip_frames(
                        clip, gif_frame_out_dir, image_path_set[1], format
                    )
                finally:
                    clip.close()
                if frame_path is not None:
                    self.file_handler.post_process(image_path_set, frame_path, delete)
            else:
                # Handle unsupported file types here
                self.event_logger.info(
                    f'[!] {lang.get_translation("skipping", self.locale)} "{self.file_handler.join_back(image_path_set)}" - {lang.get_translation("unsupported_format", self.locale)}'
                )
        # Documents can sometimes be converted to bmp, e.g. contents of docx, pdf
        for doc_path_set in file_paths[Category.DOCUMENT]:
            if doc_path_set[2] == "docx":
                office_to_frames(
                    doc_path_set=doc_path_set,
                    format=format,
                    output=output,
                    delete=delete,
                    file_handler=self.file_handler,
                    event_logger=self.event_logger,
                )
            if doc_path_set[2] == "pdf":
                pdf_path = self.file_handler.join_back(doc_path_set)
                bmp_path = os.path.abspath(
                    os.path.join(output, f"{doc_path_set[1]}.{format}")
                )
                doc = pypdf.open(pdf_path)
                bmp_frame_out_dir = _safe_out_dir(
                    os.path.join(output, doc_path_set[1])
                )
                if not os.path.exists(bmp_frame_out_dir):
                    try:
                        os.makedirs(bmp_frame_out_dir, exist_ok=True)
                    except OSError as e:
                        self.event_logger.info(
                            f"[!] {lang.get_translation('error', self.locale)}: {e} - {lang.get_translation('set_out_dir', self.locale)} {input}"
                        )
                        output = input
                        bmp_frame_out_dir = _safe_out_dir(
                            os.path.join(output, doc_path_set[1])
                        )
                for i, page_num in enumerate(range(len(doc))):
                    pix = doc.load_page(page_num).get_pixmap()
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    # Save each page as a separate BMP file
                    img.save(
                        os.path.join(
                            bmp_frame_out_dir,
                            f"{doc_path_set[1]}-{i:0{len(str(len(doc)))}}.{format}",
                        ),
                        format=format,
                    )
                doc.close()
                self.file_handler.post_process(doc_path_set, bmp_path, delete)

    def to_webp(
        self,
        input: str,
        output: str,
        file_paths: dict,
        supported_formats: dict,  # unused, aligns signature with to_frames, keep
        framerate: int,  # equally unused, same thing
        format: str,
        delete: bool,
    ) -> None:
        # Convert frames in webp format
        # Movies are converted to webps, frame by frame
        for movie_path_set in file_paths[Category.MOVIE]:
            if self.file_handler.has_visuals(movie_path_set):
                video = VideoFileClip(
                    self.file_handler.join_back(movie_path_set),
                    audio=False,
                    fps_source="tbr",
                )
                webp_frame_out_dir = _safe_out_dir(
                    os.path.join(output, movie_path_set[1])
                )
                if not os.path.exists(webp_frame_out_dir):
                    try:
                        os.makedirs(webp_frame_out_dir)
                    except OSError as e:
                        self.event_logger.info(
                            f"[!] {lang.get_translation('error', self.locale)}: {e} - {lang.get_translation('set_out_dir', self.locale)} {input}"
                        )
                        output = input
                        webp_frame_out_dir = _safe_out_dir(
                            os.path.join(output, movie_path_set[1])
                        )
                img_path = os.path.abspath(
                    os.path.join(
                        webp_frame_out_dir,
                        f"{movie_path_set[1]}-%0{len(str(int(video.duration * video.fps)))}d.{format}",
                    )
                )
                frame_names = video.write_images_sequence(
                    img_path, fps=video.fps, logger=self.prog_logger
                )
                video.close()
                if frame_names:
                    self.file_handler.post_process(
                        movie_path_set, frame_names[-1], delete
                    )
            else:
                self.event_logger.info(
                    f'[!] {lang.get_translation("skipping", self.locale)} "{self.file_handler.join_back(movie_path_set)}" - {lang.get_translation("audio_only_video", self.locale)}'
                )

        # pngs and gifs are converted to webps as well
        convertible_images = [
            ips
            for ips in file_paths[Category.IMAGE]
            if ips[2] != format
            and ips[2]
            in [
                "png",
                "jpeg",
                "jpg",
                "tiff",
                "tga",
                "eps",
                "heic",
                "heif",
                "avif",
            ]
        ]

        def _convert_image_to_webp(image_path_set: tuple):
            out_dir = self._per_file_out_dir(image_path_set, input, output)
            webp_path = self.file_handler._resolve_output_file_conflict(
                os.path.abspath(os.path.join(out_dir, f"{image_path_set[1]}.{format}"))
            )
            with Image.open(self.file_handler.join_back(image_path_set)) as img:
                img.convert("RGB").save(webp_path, format=format)
            return (image_path_set, webp_path)

        max_workers = _max_workers()
        if len(convertible_images) <= 1 or max_workers == 1:
            for image_path_set in convertible_images:
                src, webp_path = _convert_image_to_webp(image_path_set)
                self.file_handler.post_process(src, webp_path, delete)
        else:
            with ThreadPoolExecutor(max_workers=max_workers) as ex:
                futures = [
                    ex.submit(_convert_image_to_webp, ips) for ips in convertible_images
                ]
                for fut in as_completed(futures):
                    src, webp_path = fut.result()
                    self.file_handler.post_process(src, webp_path, delete)

        for image_path_set in file_paths[Category.IMAGE]:
            if image_path_set[2] == format:
                continue
            if image_path_set[2] in [
                "png",
                "jpeg",
                "jpg",
                "tiff",
                "tga",
                "eps",
                "heic",
                "heif",
                "avif",
            ]:
                continue
            elif image_path_set[2] == "gif":
                clip = VideoFileClip(self.file_handler.join_back(image_path_set))
                try:
                    gif_frame_out_dir = _safe_out_dir(
                        os.path.abspath(os.path.join(output, image_path_set[1]))
                    )
                    frame_path = _save_clip_frames(
                        clip, gif_frame_out_dir, image_path_set[1], format
                    )
                finally:
                    clip.close()
                if frame_path is not None:
                    self.file_handler.post_process(image_path_set, frame_path, delete)
            else:
                # Handle unsupported file types here
                self.event_logger.info(
                    f'[!] {lang.get_translation("skipping", self.locale)} "{self.file_handler.join_back(image_path_set)}" - {lang.get_translation("unsupported_format", self.locale)}'
                )

        # Documents may be convertable to bmps, e.g. pdfs
        for doc_path_set in file_paths[Category.DOCUMENT]:
            if doc_path_set[2] == "docx":
                office_to_frames(
                    doc_path_set=doc_path_set,
                    format=format,
                    output=output,
                    delete=delete,
                    file_handler=self.file_handler,
                    event_logger=self.event_logger,
                )
            if doc_path_set[2] == "pdf":
                pdf_path = self.file_handler.join_back(doc_path_set)
                webp_path = os.path.abspath(
                    os.path.join(output, f"{doc_path_set[1]}.{format}")
                )
                doc = pypdf.open(pdf_path)
                webp_pdf_frame_out_dir = _safe_out_dir(
                    os.path.join(output, doc_path_set[1])
                )
                if not os.path.exists(webp_pdf_frame_out_dir):
                    try:
                        os.makedirs(webp_pdf_frame_out_dir, exist_ok=True)
                    except OSError as e:
                        self.event_logger.info(
                            f"[!] {lang.get_translation('error', self.locale)}: {e} - {lang.get_translation('set_out_dir', self.locale)} {input}"
                        )
                        output = input
                        webp_pdf_frame_out_dir = _safe_out_dir(
                            os.path.join(output, doc_path_set[1])
                        )
                for i, page_num in enumerate(range(len(doc))):
                    pix = doc.load_page(page_num).get_pixmap()
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    # Save each page as a separate BMP file
                    img.save(
                        os.path.join(
                            webp_pdf_frame_out_dir,
                            f"{doc_path_set[1]}-{i:0{len(str(len(doc)))}}.{format}",
                        ),
                        format=format,
                    )
                doc.close()
                self.file_handler.post_process(doc_path_set, webp_path, delete)

    def to_heic(
        self,
        input: str,
        output: str,
        file_paths: dict,
        supported_formats: dict,  # unused, aligns signature with to_frames, so keep it
        framerate: int,
        format: str,
        delete: bool,
    ) -> None:
        # Convert still images, PDF and office documents to HEIF/HEIC/AVIF, not moviees tho
        if pillow_heif is None:
            self.event_logger.error(
                f"[!] {lang.get_translation('error', self.locale)}: "
                f"pillow-heif is required for {format} conversion. "
                f"Install it with: pip install pillow-heif"
            )
            return

        convertible_images = [
            ips
            for ips in file_paths[Category.IMAGE]
            if ips[2] not in ["heic", "heif", "avif"]
        ]

        def _convert_image(ips: tuple) -> tuple:
            out_dir = self._per_file_out_dir(ips, input, output)
            heic_path = self.file_handler._resolve_output_file_conflict(
                os.path.abspath(os.path.join(out_dir, f"{ips[1]}.{format}"))
            )
            with Image.open(self.file_handler.join_back(ips)) as img:
                if img.mode != "RGB":
                    img = img.convert("RGB")
                img.save(heic_path, format=_save_format(format))
            return (ips, heic_path)

        max_workers = _max_workers()
        if len(convertible_images) <= 1 or max_workers == 1:
            for ips in convertible_images:
                src, heic_path = _convert_image(ips)
                self.file_handler.post_process(src, heic_path, delete)
        else:
            with ThreadPoolExecutor(max_workers=max_workers) as ex:
                futures = [ex.submit(_convert_image, ips) for ips in convertible_images]
                for fut in as_completed(futures):
                    src, heic_path = fut.result()
                    self.file_handler.post_process(src, heic_path, delete)

        # Documents can be converted to HEIC, e.g. contents of docx/pdf
        for doc_path_set in file_paths[Category.DOCUMENT]:
            if doc_path_set[2] == "docx":
                office_to_frames(
                    doc_path_set=doc_path_set,
                    format=format,
                    output=output,
                    delete=delete,
                    file_handler=self.file_handler,
                    event_logger=self.event_logger,
                )
            if doc_path_set[2] == "pdf":
                pdf_path = self.file_handler.join_back(doc_path_set)
                heic_path = os.path.abspath(
                    os.path.join(output, f"{doc_path_set[1]}.{format}")
                )
                doc = pypdf.open(pdf_path)
                heic_frame_out_dir = _safe_out_dir(os.path.join(output, doc_path_set[1]))
                if not os.path.exists(heic_frame_out_dir):
                    try:
                        os.makedirs(heic_frame_out_dir, exist_ok=True)
                    except OSError as e:
                        self.event_logger.info(
                            f"[!] {lang.get_translation('error', self.locale)}: {e} - {lang.get_translation('set_out_dir', self.locale)} {input}"
                        )
                        output = input
                        heic_frame_out_dir = _safe_out_dir(
                            os.path.join(output, doc_path_set[1])
                        )
                for i, page_num in enumerate(range(len(doc))):
                    pix = doc.load_page(page_num).get_pixmap()
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    img.save(
                        os.path.join(
                            heic_frame_out_dir,
                            f"{doc_path_set[1]}-{i:0{len(str(len(doc)))}}.{format}",
                        ),
                        format=_save_format(format),
                    )
                doc.close()
                self.file_handler.post_process(doc_path_set, heic_path, delete)

    def to_gif(
        self,
        input: str,
        output: str,
        file_paths: dict,
        supported_formats: dict,  # unused, aligns signature with to_frames, so keep it
        framerate: int,
        format: str,
        delete: bool,
    ) -> None:
        # All images in the input directory are merged into one gif
        if len(file_paths[Category.IMAGE]) > 0:
            images = []
            total_images = len(file_paths[Category.IMAGE])
            processed = 0

            # Log start of GIF creation
            if hasattr(self, "prog_logger") and hasattr(
                self.prog_logger, "bars_callback"
            ):
                self.prog_logger.bars_callback("gif", "index", 0, 0)

            # Create tqdm progress bar
            progress_bar = tqdm(
                total=total_images,
                unit="img",
                leave=False,
                disable=False,
            )

            try:
                for i, image_path_set in enumerate(file_paths[Category.IMAGE], 1):
                    if image_path_set[2] == format:
                        progress_bar.update(1)  # Update for skipped images too
                        continue

                    try:
                        with Image.open(
                            self.file_handler.join_back(image_path_set)
                        ) as image:
                            images.append(image.convert("RGB"))
                        processed += 1

                        # Update progress after each image
                        if hasattr(self, "prog_logger") and hasattr(
                            self.prog_logger, "bars_callback"
                        ):
                            self.prog_logger.bars_callback("gif", "index", i, i - 1)

                        progress_bar.set_postfix(
                            {"current": os.path.basename(image_path_set[1])},
                            refresh=False,
                        )
                        progress_bar.update(1)

                    except Exception as e:
                        error_msg = f"Error processing image {os.path.basename(image_path_set[1])}: {str(e)}"
                        if hasattr(self, "prog_logger"):
                            self.prog_logger.log(error_msg)
                        progress_bar.write(error_msg)  # Show error in tqdm output
                        raise  # Re-raise to maintain original error handling

            finally:
                progress_bar.close()

            if images:
                output_path = self.file_handler._resolve_output_file_conflict(
                    os.path.abspath(os.path.join(output, f"merged.{format}"))
                )

                try:
                    images[0].save(
                        output_path,
                        save_all=True,
                        append_images=images[1:],
                    )
                    # Log completion
                    if hasattr(self, "prog_logger") and hasattr(
                        self.prog_logger, "bars_callback"
                    ):
                        self.prog_logger.bars_callback(
                            "gif", "index", total_images, total_images - 1
                        )
                except Exception as e:
                    if hasattr(self, "prog_logger"):
                        self.prog_logger.log(
                            f"Error saving GIF {output_path}: {str(e)}"
                        )
                    raise
        # Movies are converted to gifs as well, retaining 1/3 of the frames
        for i, movie_path_set in enumerate(file_paths[Category.MOVIE], 1):
            if self.file_handler.has_visuals(movie_path_set):
                try:
                    # Log start of video to GIF conversion
                    if hasattr(self, "prog_logger") and hasattr(
                        self.prog_logger, "bars_callback"
                    ):
                        self.prog_logger.bars_callback("video_gif", "index", 0, 0)

                    video = VideoFileClip(
                        self.file_handler.join_back(movie_path_set),
                        audio=False,
                        fps_source="tbr",
                    )

                    gif_path = self.file_handler._resolve_output_file_conflict(
                        os.path.abspath(
                            os.path.join(output, f"{movie_path_set[1]}.{format}")
                        )
                    )
                    target_fps = max(1, int(video.fps // 3))

                    video.write_gif(
                        gif_path,
                        fps=target_fps,
                        logger=self.prog_logger,
                    )

                    # Log completion of this video
                    if hasattr(self, "prog_logger") and hasattr(
                        self.prog_logger, "bars_callback"
                    ):
                        self.prog_logger.bars_callback("video_gif", "index", i, i - 1)

                    video.close()
                    self.file_handler.post_process(movie_path_set, gif_path, delete)
                except Exception as e:
                    if hasattr(self, "prog_logger"):
                        self.prog_logger.log(
                            f"Error converting video {movie_path_set} to GIF: {str(e)}"
                        )
                    raise  # Maintain original error handling
            else:
                self.event_logger.info(
                    f'[!] {lang.get_translation("skipping", self.locale)} "{self.file_handler.join_back(movie_path_set)}" - {lang.get_translation("audio_only_video", self.locale)}'
                )
        # Documents may be convertable to gifs, e.g. pdfs
        for doc_path_set in file_paths[Category.DOCUMENT]:
            if doc_path_set[2] == "pdf":
                pdf_path = self.file_handler.join_back(doc_path_set)
                gif_path = self.file_handler._resolve_output_file_conflict(
                    os.path.abspath(os.path.join(output, f"{doc_path_set[1]}.{format}"))
                )

                doc = pypdf.open(pdf_path)
                images = []
                for page_num in range(len(doc)):
                    pix = doc.load_page(page_num).get_pixmap()
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    images.append(img)
                if images:
                    images[0].save(
                        gif_path,
                        save_all=True,
                        append_images=images[1:],
                        duration=(len(images) * 1000 // len(images))
                        // (12 if framerate is None else framerate),
                        loop=0,
                    )
                doc.close()
                self.file_handler.post_process(doc_path_set, gif_path, delete)
            elif doc_path_set[2] in ["docx", "pptx"]:
                input_path = self.file_handler.join_back(doc_path_set)
                gif_path = self.file_handler._resolve_output_file_conflict(
                    os.path.abspath(os.path.join(output, f"{doc_path_set[1]}.{format}"))
                )

                images = []
                if doc_path_set[2] == "docx":
                    doc = docx.Document(input_path)
                    for rel in doc.part.rels.values():
                        if "image" in rel.reltype:
                            img = Image.open(BytesIO(rel.target_part.blob))
                            images.append(img.convert("RGB"))
                    frame_count = len(doc.paragraphs) or 1
                else:
                    prs = pptx.Presentation(input_path)
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if shape.shape_type == 13:  # Picture
                                image = shape.image
                                img_bytes = image.blob
                                img = Image.open(BytesIO(img_bytes))
                                images.append(img.convert("RGB"))
                    frame_count = len(prs.slides) or 1
                if images:
                    images[0].save(
                        gif_path,
                        save_all=True,
                        append_images=images[1:],
                        duration=(len(images) * 1000 // frame_count)
                        // (12 if framerate is None else framerate),
                        loop=0,
                    )
                self.file_handler.post_process(doc_path_set, gif_path, delete)

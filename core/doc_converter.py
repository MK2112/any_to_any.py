import os
import fitz
import docx
import pptx
import shutil
import mammoth
import subprocess
import utils.language_support as lang
from PIL import Image
from tqdm import tqdm
from weasyprint import HTML
from moviepy import VideoFileClip
from markdownify import markdownify
from utils.category import Category
from core.image_converter import gif_to_frames


class DocumentConverter:
    def __init__(
        self, file_handler, prog_logger, event_logger, locale: str = "English"
    ):
        self.file_handler = file_handler
        self.prog_logger = prog_logger
        self.event_logger = event_logger
        self.locale = locale

    def to_markdown(
        self, output: str, file_paths: dict, format: str, delete: bool
    ) -> None:
        # Convert Documents to Markdown
        for doc_path_set in file_paths[Category.DOCUMENT]:
            if doc_path_set[2] == "docx":
                docx_path = self.file_handler.join_back(doc_path_set)
                output_basename = doc_path_set[1]
                md_path = os.path.abspath(
                    os.path.join(output, f"{output_basename}.{format}")
                )
                image_md_dir = os.path.join(output, f"{output_basename}_images")
                os.makedirs(image_md_dir, exist_ok=True)
                image_index = 0

                # Custom image converter for mammoth
                def convert_image(image):
                    nonlocal image_index
                    extension = image.content_type.split("/")[-1]
                    image_filename = f"{output_basename}_{image_index}.{extension}"
                    image_path = os.path.join(image_md_dir, image_filename)
                    with image.open() as image_bytes:
                        buf = image_bytes.read()
                    with open(image_path, "wb") as img_file:
                        # Write the image to the file byte by byte
                        # image is of type Image
                        img_file.write(buf)
                    image_index += 1
                    # Return src attribute that points to the relative image path
                    rel_image_path = os.path.abspath(image_path)
                    return {"src": rel_image_path}

                with open(docx_path, "rb") as docx_file:
                    document = mammoth.convert_to_html(
                        docx_file,
                        convert_image=mammoth.images.img_element(convert_image),
                    )
                html_content = document.value  # Already a str, no need to encode
                # Convert HTML to Markdown
                markdown_text = markdownify(html_content)
                # Write Markdown file
                with open(md_path, "w", encoding="utf-8") as md_file:
                    md_file.write(markdown_text)
                self.file_handler.post_process(doc_path_set, md_path, delete)

    def to_pdf(self, output: str, file_paths: dict, format: str, delete: bool) -> None:
        # Convert GIFs to Frames using to_frames
        # Produces a folder with gif frame for each gif
        gif_to_frames(output, file_paths, self.file_handler)
        # Convert Images to PDF
        for image_path_set in file_paths[Category.IMAGE]:
            # Convert image to pdf
            pdf_path = ""
            if image_path_set[2] != "gif":
                doc = fitz.open()
                img = fitz.Pixmap(self.file_handler.join_back(image_path_set))
                rect = fitz.Rect(0, 0, img.width, img.height)
                page = doc.new_page(width=rect.width, height=rect.height)
                page.insert_image(rect, pixmap=img)
                pdf_path = os.path.abspath(
                    os.path.join(output, f"{image_path_set[1]}.{format}")
                )
                doc.save(pdf_path)
                doc.close()
                self.file_handler.post_process(image_path_set, pdf_path, delete)
            elif image_path_set[2] == "gif":
                # We suppose the gif was converted to frames and we have a folder of pngs
                # All pngs shall be merged into one pdf
                gif_frame_path = os.path.join(output, image_path_set[1])
                pdf_path = os.path.abspath(
                    os.path.join(output, f"{image_path_set[1]}.{format}")
                )
                doc = fitz.open()
                for frame in sorted(os.listdir(gif_frame_path)):
                    if frame.endswith(".png"):
                        img = fitz.Pixmap(os.path.join(gif_frame_path, frame))
                        rect = fitz.Rect(0, 0, img.width, img.height)
                        page = doc.new_page(width=rect.width, height=rect.height)
                        page.insert_image(rect, pixmap=img)
                doc.save(pdf_path)
                doc.close()
                # Remove the gif frame folder
                shutil.rmtree(gif_frame_path)
                self.file_handler.post_process(image_path_set, pdf_path, delete)
        # Convert Movies to PDF, because we can
        for movie_path_set in file_paths[Category.MOVIE]:
            if self.file_handler.has_visuals(movie_path_set):
                clip = VideoFileClip(
                    self.file_handler.join_back(movie_path_set),
                    audio=False,
                    fps_source="tbr",
                )
                pdf_path = os.path.abspath(
                    os.path.join(output, f"{movie_path_set[1]}.{format}")
                )
                num_digits = len(str(int(clip.duration * clip.fps)))
                doc = fitz.open()
                for i, frame in tqdm(
                    enumerate(clip.iter_frames(fps=clip.fps, dtype="uint8")),
                ):
                    frame_path = os.path.abspath(
                        os.path.join(
                            output,
                            f"{movie_path_set[1]}-temp-{i:0{num_digits}d}.png",
                        )
                    )
                    Image.fromarray(frame).save(frame_path, format="PNG")
                    img = fitz.Pixmap(frame_path)
                    rect = fitz.Rect(0, 0, img.width, img.height)
                    page = doc.new_page(width=rect.width, height=rect.height)
                    page.insert_image(rect, pixmap=img)
                    os.remove(frame_path)
                doc.save(pdf_path)
                doc.close()
                clip.close()
                self.file_handler.post_process(movie_path_set, pdf_path, delete)
        # Convert Documents to PDF
        for doc_path_set in file_paths[Category.DOCUMENT]:
            if doc_path_set[2] == "pdf":
                # If document is already a pdf, skip
                continue
            if doc_path_set[2] == "srt":
                # Convert srt to pdf
                pdf_path = os.path.abspath(
                    os.path.join(output, f"{doc_path_set[1]}.{format}")
                )
                with open(self.file_handler.join_back(doc_path_set), "r") as srt_file:
                    srt_content = srt_file.read()
                # Insert the SRT content into the PDF
                doc = fitz.open()
                page = doc.new_page()
                page.insert_text((50, 50), srt_content, fontsize=12)
                doc.save(pdf_path)
                doc.close()
                self.file_handler.post_process(doc_path_set, pdf_path, delete)
            elif doc_path_set[2] == "docx":
                pdf_path = os.path.abspath(
                    os.path.join(output, f"{doc_path_set[1]}.{format}")
                )
                docx_doc = open(self.file_handler.join_back(doc_path_set), "rb")
                # Convert docx to HTML as intermediary
                document = mammoth.convert_to_html(docx_doc)
                docx_doc.close()
                # Convert html to PDF, save that
                HTML(string=document.value.encode("utf-8")).write_pdf(pdf_path)
                self.file_handler.post_process(doc_path_set, pdf_path, delete)

    def to_subtitles(
        self, output: str, file_paths: dict, format: str, delete: bool
    ) -> None:
        # Extract Subtitles from Movies
        for movie_path_set in file_paths[Category.MOVIE]:
            input_path = self.file_handler.join_back(movie_path_set)
            out_path = os.path.abspath(os.path.join(output, f"{movie_path_set[1]}.srt"))
            self.event_logger.info(
                f"[>] {lang.get_translation('extract_subtitles', self.locale)} '{input_path}'"
            )
            try:
                # Use FFmpeg to extract subtitles
                _ = subprocess.run(
                    [
                        "ffmpeg",
                        "-i",
                        input_path,
                        "-map",
                        "0:s:0",  # Selects first subtitle stream
                        "-c:s",
                        "srt",
                        out_path,
                    ],
                    capture_output=True,
                    text=True,
                )

                if os.path.exists(out_path) and os.path.getsize(out_path) > 0:
                    self.event_logger.info(
                        f"[>] {lang.get_translation('subtitles_success', self.locale)} '{out_path}'"
                    )
                    self.file_handler.post_process(
                        movie_path_set, out_path, delete, show_status=False
                    )
                else:
                    # Try extracting closed captions when direct extract fails (found mostly in MP4 and MKV)
                    self.event_logger.info(
                        f"[!] {lang.get_translation('extract_subtitles_alt', self.locale)}"
                    )
                    _ = subprocess.run(
                        ["ffmpeg", "-i", input_path, "-c:s", format, out_path],
                        capture_output=True,
                        text=True,
                    )
                    if os.path.exists(out_path) and os.path.getsize(out_path) > 0:
                        self.event_logger.info(
                            f"[>] {lang.get_translation('embed_subtitles_success')} '{out_path}'"
                        )
                        self.file_handler.post_process(
                            movie_path_set, out_path, delete, show_status=False
                        )
                    else:
                        self.event_logger.info(
                            f"[!] {lang.get_translation('embed_subtitles_fail', self.locale)} '{input_path}'"
                        )
            except Exception as e:
                self.event_logger.info(
                    f"[!] {lang.get_translation('extract_subtitles_fail', self.locale)} {str(e)}"
                )
                try:
                    subprocess.run(["ffmpeg", "-version"], capture_output=True)
                except FileNotFoundError:
                    self.event_logger.info(
                        f"[!] {lang.get_translation('ffmpeg_not_found', self.locale)}"
                    )
                    break

    def to_office(
        self, output: str, file_paths: dict, format: str, delete: bool
    ) -> None:
        # Convert Images and Movies to Office Documents
        def _new_container():
            # Type-specific container creation
            return pptx.Presentation() if format == "pptx" else docx.Document()

        def _add_page(container):
            # Add a new page to the container, regardless of format
            return (
                container.slides.add_slide(container.slide_layouts[5])
                if format == "pptx"
                else container
            )

        def _place_img(page, img_path, full_page=False):
            # Place an image on the page, adjusting size for full page if needed
            with Image.open(img_path) as im:
                w, h = im.size
            if format == "pptx":
                max_h, max_w = container.slide_height, container.slide_width
                if full_page or h > max_h or w > max_w:
                    if h / max_h >= w / max_w:
                        page.shapes.add_picture(img_path, 0, 0, height=max_h)
                    else:
                        page.shapes.add_picture(img_path, 0, 0, width=max_w)
                else:
                    page.shapes.add_picture(img_path, 0, 0)
            elif format == "docx":
                try:
                    page.add_picture(img_path, height=docx.shared.Inches(7))
                except docx.image.exceptions.UnexpectedEndOfFileError as e:
                    self.event_logger.info(e)

        gif_to_frames(output, file_paths, self.file_handler)

        for image_path_set in file_paths[Category.IMAGE]:
            out_path = os.path.abspath(
                os.path.join(output, f"{image_path_set[1]}.{format}")
            )
            container = _new_container()
            if image_path_set[2] == "gif":
                frame_dir = os.path.join(output, image_path_set[1])
                for frame in tqdm(sorted(os.listdir(frame_dir))):
                    if not frame.endswith(".png"):
                        continue
                    page = _add_page(container)
                    _place_img(page, os.path.join(frame_dir, frame), full_page=True)
                    if format == "docx":
                        page.add_paragraph(f"Image: {frame}")
                        page.add_page_break()
                shutil.rmtree(frame_dir)
            else:
                page = _add_page(container)
                _place_img(
                    page, self.file_handler.join_back(image_path_set), full_page=True
                )
                if format == "docx":
                    page.add_paragraph(f"Image: {image_path_set[1]}")
            container.save(out_path)
            self.file_handler.post_process(image_path_set, out_path, delete)

        for movie_path_set in file_paths[Category.MOVIE]:
            if not self.file_handler.has_visuals(movie_path_set):
                continue

            out_path = os.path.abspath(
                os.path.join(output, f"{movie_path_set[1]}.{format}")
            )
            container = _new_container()
            clip = VideoFileClip(
                self.file_handler.join_back(movie_path_set),
                audio=False,
                fps_source="tbr",
            )
            digits = len(str(int(clip.duration * clip.fps)))

            for idx, frame in tqdm(
                enumerate(clip.iter_frames(fps=clip.fps, dtype="uint8"))
            ):
                frame_png = os.path.join(
                    output, f"{movie_path_set[1]}-temp-{idx:0{digits}d}.png"
                )
                Image.fromarray(frame).save(frame_png, format="PNG")
                page = _add_page(container)
                _place_img(page, frame_png, full_page=format == "pptx")
                if format == "docx":
                    page.add_paragraph(f"Frame: {idx}")
                    page.add_page_break()
                os.remove(frame_png)

            clip.close()
            container.save(out_path)
            self.file_handler.post_process(movie_path_set, out_path, delete)

        if format == "docx":
            for document_path_set in file_paths[Category.DOCUMENT]:
                if document_path_set[2] == format:
                    continue

                out_path = os.path.abspath(
                    os.path.join(output, f"{document_path_set[1]}.docx")
                )
                doc = docx.Document()

                if document_path_set[2] == "pptx":
                    pptx_path = self.file_handler.join_back(document_path_set)
                    presentation = pptx.Presentation(pptx_path)
                    for slide in tqdm(
                        presentation.slides, desc=f"{document_path_set[1]}"
                    ):
                        if slide.shapes.title.text != "":
                            # Add title to the document
                            doc.add_paragraph(
                                f"Slide {slide.slide_id} Title: {slide.shapes.title.text}"
                            )
                        for shape in slide.shapes:
                            if hasattr(shape, "image"):
                                img = shape.image
                                slide_img = os.path.join(
                                    output,
                                    f"{document_path_set[1]}-slide{slide.slide_id}.{img.ext}",
                                )
                                with open(slide_img, "wb") as f:
                                    f.write(img.blob)
                                doc.add_paragraph(f"Slide {slide.slide_id} Image:")
                                doc.add_picture(slide_img, width=docx.shared.Inches(5))
                                os.remove(slide_img)
                            if hasattr(shape, "text"):
                                if shape.text != "":
                                    doc.add_paragraph(
                                        f"Slide {slide.slide_id} Text: {shape.text}"
                                    )
                        doc.add_page_break()
                elif document_path_set[2] == "pdf":
                    pdf = fitz.open(self.file_handler.join_back(document_path_set))
                    for pnum in tqdm(range(len(pdf))):
                        page = pdf.load_page(pnum)
                        txt = page.get_text().strip()
                        if txt:
                            doc.add_paragraph(f"Page {pnum + 1} Text:")
                            for line in txt.splitlines():
                                doc.add_paragraph(line)
                        imgs = page.get_images(full=True)
                        if imgs:
                            doc.add_paragraph(f"Page {pnum + 1} Images:")
                        for i, img in enumerate(imgs):
                            xref = img[0]
                            pix = pdf.extract_image(xref)
                            img_bytes, img_ext = pix["image"], pix["ext"]
                            tmp_img = os.path.join(
                                output,
                                f"{document_path_set[1]}-page{pnum + 1}-{i}.{img_ext}",
                            )
                            with open(tmp_img, "wb") as f:
                                f.write(img_bytes)
                            doc.add_picture(tmp_img, width=docx.shared.Inches(5))
                            os.remove(tmp_img)
                        doc.add_page_break()
                doc.save(out_path)
                self.file_handler.post_process(document_path_set, out_path, delete)

    def split_pdf(self, output: str, doc_path_set: tuple, page_ranges: str, delete: bool, format: str = "pdf"):
        # output: str - output directory
        # doc_path_set: tuple - (path, name, format)
        # page_ranges: str - page ranges to split the document into
        # delete: bool - delete the original document after completion
        # format: str - format of the document
        doc_path = self.file_handler.join_back(doc_path_set)
        if format == "pdf":
            # Open the PDF document
            pdf = fitz.open(doc_path)
            total_pages = len(pdf)
            if total_pages == 0:
                pdf.close()
                return
            # Parse page_ranges into a list of tuples
            parsed_ranges = self._parse_page_ranges(page_ranges, total_pages)

            # None means any splitting efforts would be futile, skip
            if parsed_ranges is None:
                pdf.close()
                return

            # Create output directory if it doesn't exist
            os.makedirs(output, exist_ok=True)
            for i, (start, end) in enumerate(parsed_ranges):
                new_pdf = fitz.open()
                new_pdf.insert_pdf(pdf, from_page=start-1, to_page=end-1)
                # Generate output filename
                if len(parsed_ranges) == 1:
                    out_filename = f"{doc_path_set[1]}_{start}-{end}.{format}"
                else:
                    out_filename = f"{doc_path_set[1]}_split_{i+1}_{start}-{end}.{format}"
                out_path = os.path.abspath(os.path.join(output, out_filename))
                # Save new PDF
                new_pdf.save(out_path)
                new_pdf.close()
                self.event_logger.info(f"[>] {lang.get_translation('split_produced', self.locale)}: {out_path}")
            # Close source PDF
            pdf.close()
            self.file_handler.post_process(doc_path_set, out_path, delete)

    def _parse_page_ranges(self, page_ranges: str, total_pages: int) -> list[tuple[int, int]] | None:
        """
        Parse page_ranges string into list of (start, end) tuples.
        Handles: None, '1-5', '2-5,3-4', '3-6, 8-20, 23-45, rest', 'all', 'rest', 
        '3-6;15-22', '1-7;8-15;20-22;rest', '3', '35', '12-end', '2-end',
        '2-5, 3-6, 12-end', '2-end, 3-6, 12-end', '2-end, 3-6'
        """
        if not page_ranges or page_ranges.strip() == "":
            # Default to all pages
            return None #[(1, total_pages)], we don't need that, skip that
        page_ranges = page_ranges.strip()
        if page_ranges.lower() in ["all", "rest"]:
            return None #[(1, total_pages)], we don't need that, skip that

        ranges = []
        for delim in [',', ';']:
            # Split by comma or semicolon
            if delim in page_ranges:
                ranges = [r.strip() for r in page_ranges.split(delim)]
                break
        if not ranges:
            # No delimiters found
            ranges = [page_ranges.strip()]
        
        parsed_ranges = []
        # Track where 'rest' should start
        rest_start = 1
        
        for range_str in ranges:
            range_str = range_str.strip()
            if range_str.lower() == "rest":
                # 'rest': from highest processed page + 1 to end
                if parsed_ranges:
                    # Find highest end page processed so far
                    max_end = max(end for _, end in parsed_ranges)
                    rest_start = max_end + 1
                parsed_ranges.append((rest_start, total_pages))
                continue
            
            # Handle single page numbers
            if range_str.isdigit():
                page_num = int(range_str)
                if 1 <= page_num <= total_pages:
                    parsed_ranges.append((page_num, page_num))
                    rest_start = max(rest_start, page_num + 1)
                continue
            
            # Handle ranges like "1-5", "12-end", "2-end"
            if '-' in range_str:
                parts = range_str.split('-', 1)
                if len(parts) == 2:
                    start_str, end_str = parts[0].strip(), parts[1].strip()
                    
                    # Parse start
                    try:
                        start = int(start_str)
                    except ValueError:
                        continue # Some invalid range
                    
                    # Parse end
                    if end_str.lower() == "end":
                        end = total_pages
                    else:
                        try:
                            end = int(end_str)
                        except ValueError:
                            continue # Some invalid range
                    
                    # Validate, add range
                    if 1 <= start <= total_pages and 1 <= end <= total_pages and start <= end:
                        parsed_ranges.append((start, end))
                        rest_start = max(rest_start, end + 1)
        
        # Remove duplicate range tuples, preserve order
        # Ranges still are allowed to overlap, must each be unique though
        seen = set()
        unique_ranges = []
        for range_tuple in parsed_ranges:
            if range_tuple not in seen:
                seen.add(range_tuple)
                unique_ranges.append(range_tuple)
        
        if not unique_ranges:
            # If no valid ranges parsed, default to all pages, which is None
            unique_ranges = None
        
        return unique_ranges
import os
import docx
import pptx
import fitz
import PyPDF2
import utils.language_support as lang
from PIL import Image
from tqdm import tqdm
from moviepy import VideoFileClip
from utils.category import Category


def office_to_frames(
    doc_path_set: tuple,
    format: str,
    output: str,
    delete: bool,
    file_handler,
    event_logger,
) -> None:
    # Utility function available beyond DocumentConverter scope
    # Used e.g. in MovieConverter, placing it here makes it accessible easily
    full_path = file_handler.join_back(doc_path_set)
    output_dir = os.path.join(output, doc_path_set[1])
    os.makedirs(output_dir, exist_ok=True)
    try:
        if doc_path_set[2] == "docx":
            doc = docx.Document(full_path)
            for i, rel in tqdm(
                enumerate(doc.part.rels.values()), desc=f"{doc_path_set[1]}"
            ):
                if "image" in rel.reltype:
                    img_bytes = rel.target_part.blob
                    out_path = os.path.join(
                        output_dir, f"{doc_path_set[1]}-{i}.{format}"
                    )
                    with open(out_path, "wb") as f:
                        f.write(img_bytes)
            file_handler.post_process(doc_path_set, out_path, delete)
        elif doc_path_set[2] == "pptx":
            prs = pptx.Presentation(full_path)
            img_count = 0
            for i, slide in tqdm(enumerate(prs.slides), desc=f"{doc_path_set[1]}"):
                for shape in slide.shapes:
                    if shape.shape_type == 13:  # Picture type
                        image = shape.image
                        img_bytes = image.blob
                        out_path = os.path.join(
                            output_dir, f"{doc_path_set[1]}-{img_count}.{format}"
                        )
                        with open(out_path, "wb") as f:
                            f.write(img_bytes)
                        img_count += 1
            file_handler.post_process(doc_path_set, out_path, delete)
    except Exception as e:
        event_logger.error(e)


def gif_to_frames(output: str, file_paths: dict, file_handler) -> None:
    # Convert GIFs to frames, place those in a folder
    gifs = [
        image_path
        for image_path in file_paths[Category.IMAGE]
        if image_path[2] == "gif"
    ]
    if len(gifs) > 0:
        for image_path_set in gifs:
            clip = VideoFileClip(file_handler.join_back(image_path_set), audio=False)
            # Calculate zero-padding width based on total number of frames
            total_frames = int(clip.duration * clip.fps)
            num_digits = len(str(total_frames))
            # Create a dedicated folder for the gif to store its frames
            if not os.path.exists(os.path.join(output, image_path_set[1])):
                os.makedirs(os.path.join(output, image_path_set[1]), exist_ok=True)
            for i, frame in enumerate(clip.iter_frames(fps=clip.fps, dtype="uint8")):
                frame_filename = f"{image_path_set[1]}-{i:0{num_digits}d}.png"
                frame_path = os.path.abspath(
                    os.path.join(output, image_path_set[1], frame_filename)
                )
                Image.fromarray(frame).save(frame_path, format="PNG")
            clip.close()


class ImageConverter:
    def __init__(
        self, file_handler, prog_logger, event_logger, locale: str = "English"
    ):
        self.file_handler = file_handler
        self.prog_logger = prog_logger
        self.event_logger = event_logger
        self.locale = locale

    def to_frames(
        self,
        input: str,
        output: str,
        file_paths: dict,
        supported_formats: dict,  # self._supported_formats
        framerate: int,
        format: str,
        delete: bool,
    ) -> None:
        # Converting to image frame sets
        # This works for images and movies only
        format = "jpeg" if format == "jpg" else format

        # Handle GIF files first
        gif_to_frames(output, file_paths, self.file_handler)

        # Process regular image files
        for image_path_set in file_paths[Category.IMAGE]:
            try:
                if image_path_set[2] == format:
                    continue

                if image_path_set[2] == "gif":
                    # gif_to_frames already processed these
                    self.file_handler.post_process(image_path_set, output, delete)
                    continue

                # Ensure output directory exists
                output_dir = os.path.dirname(os.path.join(output, image_path_set[1]))
                os.makedirs(output_dir, exist_ok=True)

                # Construct full output path
                img_path = os.path.abspath(
                    os.path.join(output, f"{image_path_set[1]}.{format}")
                )

                # Convert and save the image
                with Image.open(self.file_handler.join_back(image_path_set)) as img:
                    # Convert to RGB if needed (important for WebP to PNG)
                    if img.mode != "RGB":
                        img = img.convert("RGB")
                    img.save(img_path, format=format.upper())

                # Verify the file was created before deleting source
                if os.path.exists(img_path):
                    self.file_handler.post_process(image_path_set, img_path, delete)
                else:
                    self.event_logger.error(
                        f"[!] {lang.get_translation('conversion_failed', self.locale)}: "
                        f"{self.file_handler.join_back(image_path_set)}"
                    )

            except Exception as e:
                self.event_logger.error(
                    f"[!] {lang.get_translation('error', self.locale)} "
                    f"converting {self.file_handler.join_back(image_path_set)}: {str(e)}"
                )
                # Don't delete source file if conversion failed
                delete = False
                continue
        # Convert documents to image frames
        for doc_path_set in file_paths[Category.DOCUMENT]:
            if doc_path_set[2] == format:
                continue
            if not os.path.exists(os.path.join(output, doc_path_set[1])):
                try:
                    os.makedirs(os.path.join(output, doc_path_set[1]), exist_ok=True)
                except OSError as e:
                    self.event_logger.info(
                        f"[!] {lang.get_translation('error', self.locale)}: {e} - {lang.get_translation('set_out_dir', self.locale)} {input}"
                    )
                    output = input
            if doc_path_set[2] in ["docx", "pptx"]:
                # Read all images from docx, write to os.path.join(self.output, doc_path_set[1])
                office_to_frames(
                    doc_path_set=doc_path_set,
                    format=format,
                    output=output,
                    delete=delete,
                    file_handler=self.file_handler,
                    event_logger=self.event_logger,
                )
            if doc_path_set[2] == "pdf":
                # Per page, convert pdf to image
                pdf_path = self.file_handler.join_back(doc_path_set)
                pdf = PyPDF2.PdfReader(pdf_path)
                img_path = os.path.abspath(
                    os.path.join(
                        os.path.join(output, doc_path_set[1]),
                        f"{doc_path_set[1]}-%{len(str(len(pdf.pages)))}d.{format}",
                    )
                )

                try:
                    if not os.path.exists(os.path.dirname(img_path)):
                        os.makedirs(os.path.dirname(img_path), exist_ok=True)
                    output = os.path.dirname(img_path)
                except OSError as e:
                    self.event_logger.info(
                        f"[!] {lang.get_translation('error', self.locale)}: {e} - {lang.get_translation('set_out_dir', self.locale)} {input}"
                    )
                    output = input

                pdf_document = fitz.open(pdf_path)

                for page_num in range(len(pdf_document)):
                    pix = pdf_document[page_num].get_pixmap()
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    img_file = img_path % (page_num + 1)
                    img.save(img_file, format.upper())

                pdf_document.close()
                self.file_handler.post_process(doc_path_set, img_path, delete)

        # Audio cant be image-framed, movies certrainly can
        for movie_path_set in file_paths[Category.MOVIE]:
            if movie_path_set[2] not in supported_formats[Category.MOVIE]:
                self.event_logger.info(
                    f"[!] {lang.get_translation('movie_format_unsupported', self.locale)} {movie_path_set[2]} - {lang.get_translation('skipping', self.locale)}"
                )
                continue
            if self.file_handler.has_visuals(movie_path_set):
                video = VideoFileClip(
                    self.file_handler.join_back(movie_path_set),
                    audio=False,
                    fps_source="tbr",
                )
                try:
                    if not os.path.exists(os.path.join(output, movie_path_set[1])):
                        os.makedirs(
                            os.path.join(output, movie_path_set[1]), exist_ok=True
                        )
                except OSError as e:
                    self.event_logger.info(
                        f"[!] {lang.get_translation('error', self.locale)}: {e} - {lang.get_translation('set_out_dir', self.locale)} {input}"
                    )
                    output = input
                img_path = os.path.abspath(
                    os.path.join(
                        os.path.join(output, movie_path_set[1]),
                        f"{movie_path_set[1]}-%{len(str(int(video.duration * video.fps)))}d.{format}",
                    )
                )

                video.write_images_sequence(
                    img_path, fps=video.fps, logger=self.prog_logger
                )
                video.close()
                self.file_handler.post_process(movie_path_set, img_path, delete)
            else:
                self.event_logger.info(
                    f'[!] {lang.get_translation("skipping", self.locale)} "{self.file_handler.join_back(movie_path_set)}" - {lang.get_translation("audio_only_video", self.locale)}'
                )
                self.file_handler.post_process(movie_path_set, img_path, delete)

    def to_bmp(
        self,
        input: str,
        output: str,
        file_paths: dict,
        supported_formats: dict,  # unused, aligns signature with to_frames, so keep it
        framerate: int,
        format: str,
        delete: bool,
    ) -> None:
        for movie_path_set in file_paths[Category.MOVIE]:
            # Movies are converted to bmps frame by frame
            if self.file_handler.has_visuals(movie_path_set):
                video = VideoFileClip(
                    self.file_handler.join_back(movie_path_set),
                    audio=False,
                    fps_source="tbr",
                )
                bmp_path = os.path.join(output, f"{movie_path_set[1]}.{format}")
                # Split video into individual bmp frame images at original framerate
                for _, frame in enumerate(
                    video.iter_frames(fps=video.fps, dtype="uint8")
                ):
                    frame.save(
                        f"{bmp_path}-%{len(str(int(video.duration * video.fps)))}d.{format}",
                        format=format,
                    )
                self.file_handler.post_process(movie_path_set, bmp_path, delete)
            else:
                self.event_logger.info(
                    f'[!] {lang.get_translation("skipping", self.locale)} "{self.file_handler.join_back(movie_path_set)}" - {lang.get_translation("audio_only_video", self.locale)}'
                )
        for image_path_set in file_paths[Category.IMAGE]:
            # Pngs and gifs are converted to bmps as well
            if image_path_set[2] == format:
                continue
            if image_path_set[2] in ["png", "jpeg", "jpg", "tiff", "tga", "eps"]:
                bmp_path = os.path.join(output, f"{image_path_set[1]}.{format}")
                with Image.open(self.file_handler.join_back(image_path_set)) as img:
                    img.convert("RGB").save(bmp_path, format=format)
                self.file_handler.post_process(image_path_set, bmp_path, delete)
            elif image_path_set[2] == "gif":
                clip = VideoFileClip(self.file_handler.join_back(image_path_set))
                for _, frame in enumerate(
                    clip.iter_frames(fps=clip.fps, dtype="uint8")
                ):
                    frame_path = os.path.join(
                        output,
                        f"{image_path_set[1]}-%{len(str(int(clip.duration * clip.fps)))}d.{format}",
                    )
                    Image.fromarray(frame).convert("RGB").save(
                        frame_path, format=format
                    )
                self.file_handler.post_process(image_path_set, frame_path, delete)
            else:
                # Handle unsupported file types here
                self.event_logger.info(
                    f'[!] {lang.get_translation("skipping", self.locale)} "{self.file_handler.join_back(image_path_set)}" - {lang.get_translation("unsupported_format", self.locale)}'
                )
        # Documents can sometimes be converted to bmp, e.g. contents of docx, pdf
        for doc_path_set in file_paths[Category.DOCUMENT]:
            if doc_path_set[2] == "docx":
                office_to_frames(
                    doc_path_set=doc_path_set,
                    format=format,
                    output=output,
                    delete=delete,
                    file_handler=self.file_handler,
                    event_logger=self.event_logger,
                )
            if doc_path_set[2] == "pdf":
                pdf_path = self.file_handler.join_back(doc_path_set)
                bmp_path = os.path.abspath(
                    os.path.join(output, f"{doc_path_set[1]}.{format}")
                )
                doc = fitz.open(pdf_path)
                if not os.path.exists(os.path.join(output, doc_path_set[1])):
                    try:
                        os.makedirs(
                            os.path.join(output, doc_path_set[1]), exist_ok=True
                        )
                    except OSError as e:
                        self.event_logger.info(
                            f"[!] {lang.get_translation('error', self.locale)}: {e} - {lang.get_translation('set_out_dir', self.locale)} {input}"
                        )
                        output = input
                for i, page_num in enumerate(range(len(doc))):
                    pix = doc.load_page(page_num).get_pixmap()
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    img = img.convert("RGB")
                    # Save each page as a separate BMP file
                    img.save(
                        os.path.join(
                            output,
                            doc_path_set[1],
                            f"{doc_path_set[1]}-{i:0{len(str(len(doc)))}}.{format}",
                        ),
                        format=format,
                    )
                doc.close()
                self.file_handler.post_process(doc_path_set, bmp_path, delete)

    def to_webp(
        self,
        input: str,
        output: str,
        file_paths: dict,
        supported_formats: dict,  # unused, aligns signature with to_frames, so keep it
        framerate: int,
        format: str,
        delete: bool,
    ) -> None:
        # Convert frames in webp format
        # Movies are converted to webps, frame by frame
        for movie_path_set in file_paths[Category.MOVIE]:
            if self.file_handler.has_visuals(movie_path_set):
                video = VideoFileClip(
                    self.file_handler.join_back(movie_path_set),
                    audio=False,
                    fps_source="tbr",
                )
                if not os.path.exists(os.path.join(output, movie_path_set[1])):
                    try:
                        os.makedirs(os.path.join(output, movie_path_set[1]))
                    except OSError as e:
                        self.event_logger.info(
                            f"[!] {lang.get_translation('error', self.locale)}: {e} - {lang.get_translation('set_out_dir', self.locale)} {input}"
                        )
                        output = input
                img_path = os.path.abspath(
                    os.path.join(
                        os.path.join(output, movie_path_set[1]),
                        f"{movie_path_set[1]}-%{len(str(int(video.duration * video.fps)))}d.{format}",
                    )
                )
                video.write_images_sequence(
                    img_path, fps=video.fps, logger=self.prog_logger
                )
                video.close()
                self.file_handler.post_process(movie_path_set, img_path, delete)
            else:
                self.event_logger.info(
                    f'[!] {lang.get_translation("skipping", self.locale)} "{self.file_handler.join_back(movie_path_set)}" - {lang.get_translation("audio_only_video", self.locale)}'
                )

        # pngs and gifs are converted to webps as well
        for image_path_set in file_paths[Category.IMAGE]:
            if image_path_set[2] == format:
                continue
            if image_path_set[2] in ["png", "jpeg", "jpg", "tiff", "tga", "eps"]:
                webp_path = os.path.join(output, f"{image_path_set[1]}.{format}")
                with Image.open(self.file_handler.join_back(image_path_set)) as img:
                    img.convert("RGB").save(webp_path, format=format)
                self.file_handler.post_process(image_path_set, webp_path, delete)
            elif image_path_set[2] == "gif":
                clip = VideoFileClip(self.file_handler.join_back(image_path_set))
                for _, frame in enumerate(
                    clip.iter_frames(fps=clip.fps, dtype="uint8")
                ):
                    frame_path = os.path.join(
                        output,
                        f"{image_path_set[1]}-%{len(str(int(clip.duration * clip.fps)))}d.{format}",
                    )
                    Image.fromarray(frame).convert("RGB").save(
                        frame_path, format=format
                    )
                self.file_handler.post_process(image_path_set, frame_path, delete)
            else:
                # Handle unsupported file types here
                self.event_logger.info(
                    f'[!] {lang.get_translation("skipping", self.locale)} "{self.file_handler.join_back(image_path_set)}" - {lang.get_translation("unsupported_format", self.locale)}'
                )

        # Documents may be convertable to bmps, e.g. pdfs
        for doc_path_set in file_paths[Category.DOCUMENT]:
            if doc_path_set[2] == "docx":
                office_to_frames(
                    doc_path_set=doc_path_set,
                    format=format,
                    output=output,
                    delete=delete,
                    file_handler=self.file_handler,
                    event_logger=self.event_logger,
                )
            if doc_path_set[2] == "pdf":
                pdf_path = self.file_handler.join_back(doc_path_set)
                bmp_path = os.path.abspath(
                    os.path.join(output, f"{doc_path_set[1]}.{format}")
                )
                doc = fitz.open(pdf_path)
                if not os.path.exists(os.path.join(output, doc_path_set[1])):
                    try:
                        os.makedirs(
                            os.path.join(output, doc_path_set[1]), exist_ok=True
                        )
                    except OSError as e:
                        self.event_logger.info(
                            f"[!] {lang.get_translation('error', self.locale)}: {e} - {lang.get_translation('set_out_dir', self.locale)} {input}"
                        )
                        output = input
                for i, page_num in enumerate(range(len(doc))):
                    pix = doc.load_page(page_num).get_pixmap()
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    img = img.convert("RGB")
                    # Save each page as a separate BMP file
                    img.save(
                        os.path.join(
                            output,
                            doc_path_set[1],
                            f"{doc_path_set[1]}-{i:0{len(str(len(doc)))}}.{format}",
                        ),
                        format=format,
                    )
                doc.close()
                self.file_handler.post_process(doc_path_set, bmp_path, delete)

    def to_gif(
        self,
        input: str,
        output: str,
        file_paths: dict,
        supported_formats: dict,  # unused, aligns signature with to_frames, so keep it
        framerate: int,
        format: str,
        delete: bool,
    ) -> None:
        # All images in the input directory are merged into one gif
        if len(file_paths[Category.IMAGE]) > 0:
            images = []
            total_images = len(file_paths[Category.IMAGE])
            processed = 0

            # Log start of GIF creation
            if hasattr(self, "prog_logger") and hasattr(
                self.prog_logger, "bars_callback"
            ):
                self.prog_logger.bars_callback("gif", "index", 0, 0)

            # Create tqdm progress bar
            progress_bar = tqdm(
                total=total_images,
                unit="img",
                leave=False,
                disable=not getattr(self, "verbose", True),  # Respect verbosity setting
            )

            try:
                for i, image_path_set in enumerate(file_paths[Category.IMAGE], 1):
                    if image_path_set[2] == format:
                        progress_bar.update(1)  # Update for skipped images too
                        continue

                    try:
                        with Image.open(
                            self.file_handler.join_back(image_path_set)
                        ) as image:
                            images.append(image.convert("RGB"))
                        processed += 1

                        # Update progress after each image
                        if hasattr(self, "prog_logger") and hasattr(
                            self.prog_logger, "bars_callback"
                        ):
                            self.prog_logger.bars_callback("gif", "index", i, i - 1)

                        progress_bar.set_postfix(
                            {"current": os.path.basename(image_path_set[1])},
                            refresh=False,
                        )
                        progress_bar.update(1)

                    except Exception as e:
                        error_msg = f"Error processing image {os.path.basename(image_path_set[1])}: {str(e)}"
                        if hasattr(self, "prog_logger"):
                            self.prog_logger.log(error_msg)
                        progress_bar.write(error_msg)  # Show error in tqdm output
                        raise  # Re-raise to maintain original error handling

            finally:
                progress_bar.close()

            if images:
                output_path = os.path.join(output, f"{output}_merged.{format}")
                try:
                    images[0].save(
                        output_path,
                        save_all=True,
                        append_images=images[1:],
                    )
                    # Log completion
                    if hasattr(self, "prog_logger") and hasattr(
                        self.prog_logger, "bars_callback"
                    ):
                        self.prog_logger.bars_callback(
                            "gif", "index", total_images, total_images - 1
                        )
                except Exception as e:
                    if hasattr(self, "prog_logger"):
                        self.prog_logger.log(
                            f"Error saving GIF {output_path}: {str(e)}"
                        )
                    raise
        # Movies are converted to gifs as well, retaining 1/3 of the frames
        for i, movie_path_set in enumerate(file_paths[Category.MOVIE], 1):
            if self.file_handler.has_visuals(movie_path_set):
                try:
                    # Log start of video to GIF conversion
                    if hasattr(self, "prog_logger") and hasattr(
                        self.prog_logger, "bars_callback"
                    ):
                        self.prog_logger.bars_callback("video_gif", "index", 0, 0)

                    video = VideoFileClip(
                        self.file_handler.join_back(movie_path_set),
                        audio=False,
                        fps_source="tbr",
                    )
                    gif_path = os.path.join(output, f"{movie_path_set[1]}.{format}")

                    # Calculate target fps, ensuring it's at least 1
                    target_fps = max(1, int(video.fps // 3))

                    # Write GIF with progress logging
                    video.write_gif(
                        gif_path,
                        fps=target_fps,
                        logger=self.prog_logger,
                        verbose=False,  # Disable moviepy's default progress bar
                    )

                    # Log completion of this video
                    if hasattr(self, "prog_logger") and hasattr(
                        self.prog_logger, "bars_callback"
                    ):
                        self.prog_logger.bars_callback("video_gif", "index", i, i - 1)

                except Exception as e:
                    if hasattr(self, "prog_logger"):
                        self.prog_logger.log(
                            f"Error converting video {movie_path_set} to GIF: {str(e)}"
                        )
                    raise  # Re-raise to maintain original error handling
                video.close()
                self.file_handler.post_process(movie_path_set, gif_path, delete)
            else:
                self.event_logger.info(
                    f'[!] {lang.get_translation("skipping", self.locale)} "{self.file_handler.join_back(movie_path_set)}" - {lang.get_translation("audio_only_video", self.locale)}'
                )
        # Documents may be convertable to gifs, e.g. pdfs
        for doc_path_set in file_paths[Category.DOCUMENT]:
            if doc_path_set[2] == "pdf":
                pdf_path = self.file_handler.join_back(doc_path_set)
                gif_path = os.path.abspath(
                    os.path.join(output, f"{doc_path_set[1]}.{format}")
                )
                doc = fitz.open(pdf_path)
                images = []
                for page_num in range(len(doc)):
                    pix = doc.load_page(page_num).get_pixmap()
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    images.append(img.convert("RGB"))
                if images:
                    images[0].save(
                        gif_path,
                        save_all=True,
                        append_images=images[1:],
                        duration=(len(images) * 1000 // len(doc))
                        // (12 if framerate is None else framerate),
                        loop=0,
                    )
                doc.close()
                self.file_handler.post_process(doc_path_set, gif_path, delete)
            elif doc_path_set[2] in ["docx", "pptx"]:
                input_path = self.file_handler.join_back(doc_path_set)
                gif_path = os.path.abspath(
                    os.path.join(output, f"{doc_path_set[1]}.{format}")
                )
                images = []
                if doc_path_set[2] == "docx":
                    doc = docx.Document(input_path)
                    for rel in doc.part.rels.values():
                        if "image" in rel.reltype:
                            img = Image.open(rel.target_part.blob)
                            images.append(img.convert("RGB"))
                    frame_count = len(doc.paragraphs) or 1
                else:
                    prs = pptx.Presentation(input_path)
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if shape.shape_type == 13:  # Picture
                                image = shape.image
                                img_bytes = image.blob
                                img = Image.open(img_bytes)
                                images.append(img.convert("RGB"))
                    frame_count = len(prs.slides) or 1
                if images:
                    images[0].save(
                        gif_path,
                        save_all=True,
                        append_images=images[1:],
                        duration=(len(images) * 1000 // frame_count)
                        // (12 if framerate is None else framerate),
                        loop=0,
                    )
                self.file_handler.post_process(doc_path_set, gif_path, delete)
